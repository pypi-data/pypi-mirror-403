"""
Document Loaders
================

Loaders for various document formats and sources.
"""

from abc import ABC, abstractmethod
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
import os


class Document:
    """A document with content and metadata."""
    
    def __init__(self, content: str, metadata: Optional[Dict[str, Any]] = None):
        self.content = content
        self.metadata = metadata or {}
    
    def __repr__(self) -> str:
        return f"Document(len={len(self.content)}, metadata={list(self.metadata.keys())})"


class BaseLoader(ABC):
    """Base class for document loaders."""
    
    @abstractmethod
    def load(self) -> List[Document]:
        """Load documents."""
        pass
    
    def load_and_split(self, splitter: Any = None) -> List[Document]:
        """Load and optionally split documents."""
        docs = self.load()
        if splitter:
            return splitter.split_documents(docs)
        return docs


# =============================================================================
# File Loaders
# =============================================================================

class TextLoader(BaseLoader):
    """Load plain text files with automatic encoding detection."""
    
    def __init__(
        self,
        path: str,
        encoding: Optional[str] = None,
        autodetect_encoding: bool = True,
    ):
        self.path = path
        self.encoding = encoding
        self.autodetect_encoding = autodetect_encoding
    
    def _detect_encoding(self, file_path: str) -> str:
        """Detect file encoding using chardet or fallback to utf-8."""
        try:
            import chardet
            with open(file_path, "rb") as f:
                raw_data = f.read(10000)  # Read first 10KB
            result = chardet.detect(raw_data)
            return result.get("encoding") or "utf-8"
        except ImportError:
            return "utf-8"
    
    def load(self) -> List[Document]:
        if self.encoding:
            encoding = self.encoding
        elif self.autodetect_encoding:
            encoding = self._detect_encoding(self.path)
        else:
            encoding = "utf-8"
        
        try:
            with open(self.path, "r", encoding=encoding) as f:
                content = f.read()
        except UnicodeDecodeError:
            # Fallback: try with errors='replace'
            with open(self.path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
        
        return [Document(content, {
            "source": self.path,
            "encoding": encoding,
        })]


class PDFLoader(BaseLoader):
    """Load PDF files with automatic fallback to pdfplumber."""
    
    def __init__(
        self,
        path: str,
        extract_images: bool = False,
        use_pdfplumber_fallback: bool = True,
    ):
        self.path = path
        self.extract_images = extract_images
        self.use_pdfplumber_fallback = use_pdfplumber_fallback
    
    def _try_pypdf(self) -> Optional[List[Document]]:
        """Try to extract text using pypdf."""
        try:
            import pypdf
            reader = pypdf.PdfReader(self.path)
            pages = []
            total_text = 0
            
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                total_text += len(text)
                pages.append(Document(text, {
                    "source": self.path,
                    "page": i,
                    "total_pages": len(reader.pages),
                    "extractor": "pypdf",
                }))
            
            # If we got very little text, return None to trigger fallback
            if total_text < 100 and self.use_pdfplumber_fallback:
                return None
            
            return pages
        except ImportError:
            return None
        except Exception:
            return None
    
    def _try_pdfplumber(self) -> Optional[List[Document]]:
        """Try to extract text using pdfplumber (better for tables)."""
        try:
            import pdfplumber
            pages = []
            
            with pdfplumber.open(self.path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    
                    # Also extract tables if present
                    tables = page.extract_tables()
                    if tables:
                        for table in tables:
                            text += "\n\n[TABLE]\n"
                            for row in table:
                                text += " | ".join([str(c) if c else "" for c in row]) + "\n"
                    
                    pages.append(Document(text, {
                        "source": self.path,
                        "page": i,
                        "total_pages": len(pdf.pages),
                        "extractor": "pdfplumber",
                    }))
            
            return pages
        except ImportError:
            return None
        except Exception:
            return None
    
    def load(self) -> List[Document]:
        # Try pypdf first
        result = self._try_pypdf()
        if result is not None:
            return result
        
        # Fallback to pdfplumber
        if self.use_pdfplumber_fallback:
            result = self._try_pdfplumber()
            if result is not None:
                return result
        
        # If both fail, raise helpful error
        raise ImportError(
            "PDF extraction failed. Install one of:\n"
            "  pip install pypdf\n"
            "  pip install pdfplumber"
        )


class PDFPlumberLoader(BaseLoader):
    """Load PDF files using pdfplumber (better for tables)."""
    
    def __init__(self, path: str):
        self.path = path
    
    def load(self) -> List[Document]:
        try:
            import pdfplumber
            pages = []
            with pdfplumber.open(self.path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    pages.append(Document(text, {"source": self.path, "page": i}))
            return pages
        except ImportError:
            raise ImportError("pdfplumber required. pip install pdfplumber")


class DOCXLoader(BaseLoader):
    """Load DOCX files."""
    
    def __init__(self, path: str):
        self.path = path
    
    def load(self) -> List[Document]:
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(self.path)
            content = "\n".join([para.text for para in doc.paragraphs])
            return [Document(content, {"source": self.path})]
        except ImportError:
            raise ImportError("python-docx required. pip install python-docx")


class CSVLoader(BaseLoader):
    """Load CSV files."""
    
    def __init__(
        self,
        path: str,
        source_column: Optional[str] = None,
        content_columns: Optional[List[str]] = None,
    ):
        self.path = path
        self.source_column = source_column
        self.content_columns = content_columns
    
    def load(self) -> List[Document]:
        import csv
        docs = []
        with open(self.path, "r", encoding="utf-8") as f:
            reader = csv.DictReader(f)
            for i, row in enumerate(reader):
                if self.content_columns:
                    content = " ".join([str(row.get(c, "")) for c in self.content_columns])
                else:
                    content = " ".join([str(v) for v in row.values()])
                
                metadata = {"source": self.path, "row": i}
                if self.source_column and self.source_column in row:
                    metadata["source"] = row[self.source_column]
                
                docs.append(Document(content, metadata))
        return docs


class JSONLoader(BaseLoader):
    """Load JSON files."""
    
    def __init__(
        self,
        path: str,
        jq_schema: Optional[str] = None,
        content_key: Optional[str] = None,
    ):
        self.path = path
        self.jq_schema = jq_schema
        self.content_key = content_key
    
    def load(self) -> List[Document]:
        import json
        with open(self.path, "r", encoding="utf-8") as f:
            data = json.load(f)
        
        if isinstance(data, list):
            docs = []
            for i, item in enumerate(data):
                if self.content_key and isinstance(item, dict):
                    content = str(item.get(self.content_key, item))
                else:
                    content = json.dumps(item) if isinstance(item, dict) else str(item)
                docs.append(Document(content, {"source": self.path, "index": i}))
            return docs
        else:
            content = json.dumps(data) if isinstance(data, dict) else str(data)
            return [Document(content, {"source": self.path})]


class JSONLLoader(BaseLoader):
    """Load JSONL (JSON Lines) files."""
    
    def __init__(self, path: str, content_key: Optional[str] = None):
        self.path = path
        self.content_key = content_key
    
    def load(self) -> List[Document]:
        import json
        docs = []
        with open(self.path, "r", encoding="utf-8") as f:
            for i, line in enumerate(f):
                if line.strip():
                    item = json.loads(line)
                    if self.content_key and isinstance(item, dict):
                        content = str(item.get(self.content_key, item))
                    else:
                        content = json.dumps(item)
                    docs.append(Document(content, {"source": self.path, "line": i}))
        return docs


class MarkdownLoader(BaseLoader):
    """Load Markdown files."""
    
    def __init__(self, path: str):
        self.path = path
    
    def load(self) -> List[Document]:
        with open(self.path, "r", encoding="utf-8") as f:
            content = f.read()
        return [Document(content, {"source": self.path, "format": "markdown"})]


class HTMLLoader(BaseLoader):
    """Load HTML files."""
    
    def __init__(self, path: str):
        self.path = path
    
    def load(self) -> List[Document]:
        try:
            from bs4 import BeautifulSoup
            with open(self.path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text(separator="\n", strip=True)
            return [Document(text, {"source": self.path, "format": "html"})]
        except ImportError:
            raise ImportError("beautifulsoup4 required. pip install beautifulsoup4")


class XMLLoader(BaseLoader):
    """Load XML files."""
    
    def __init__(self, path: str):
        self.path = path
    
    def load(self) -> List[Document]:
        try:
            from lxml import etree
            tree = etree.parse(self.path)
            content = etree.tostring(tree, encoding="unicode", method="text")
            return [Document(content, {"source": self.path, "format": "xml"})]
        except ImportError:
            import xml.etree.ElementTree as ET
            tree = ET.parse(self.path)
            root = tree.getroot()
            content = ET.tostring(root, encoding="unicode", method="text")
            return [Document(content, {"source": self.path, "format": "xml"})]


class ExcelLoader(BaseLoader):
    """Load Excel files."""
    
    def __init__(self, path: str, sheet_name: Optional[str] = None):
        self.path = path
        self.sheet_name = sheet_name
    
    def load(self) -> List[Document]:
        try:
            import pandas as pd
            df = pd.read_excel(self.path, sheet_name=self.sheet_name)
            content = df.to_string()
            return [Document(content, {"source": self.path, "format": "excel"})]
        except ImportError:
            raise ImportError("pandas and openpyxl required. pip install pandas openpyxl")


class PowerPointLoader(BaseLoader):
    """Load PowerPoint files."""
    
    def __init__(self, path: str):
        self.path = path
    
    def load(self) -> List[Document]:
        try:
            from pptx import Presentation
            prs = Presentation(self.path)
            slides = []
            for i, slide in enumerate(prs.slides):
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
                slides.append(Document("\n".join(text_parts), {"source": self.path, "slide": i}))
            return slides
        except ImportError:
            raise ImportError("python-pptx required. pip install python-pptx")


class CodeLoader(BaseLoader):
    """Load source code files."""
    
    LANGUAGE_EXTENSIONS = {
        ".py": "python",
        ".js": "javascript",
        ".ts": "typescript",
        ".java": "java",
        ".cpp": "cpp",
        ".c": "c",
        ".go": "go",
        ".rs": "rust",
        ".rb": "ruby",
        ".php": "php",
    }
    
    def __init__(self, path: str, language: Optional[str] = None):
        self.path = path
        self.language = language
    
    def load(self) -> List[Document]:
        ext = Path(self.path).suffix.lower()
        language = self.language or self.LANGUAGE_EXTENSIONS.get(ext, "unknown")
        
        with open(self.path, "r", encoding="utf-8") as f:
            content = f.read()
        
        return [Document(content, {"source": self.path, "language": language})]


class EmailLoader(BaseLoader):
    """Load email files (.eml)."""
    
    def __init__(self, path: str):
        self.path = path
    
    def load(self) -> List[Document]:
        import email
        from email.policy import default
        
        with open(self.path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=default)
        
        subject = msg.get("Subject", "")
        sender = msg.get("From", "")
        body = ""
        
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body = part.get_content()
                    break
        else:
            body = msg.get_content()
        
        content = f"Subject: {subject}\nFrom: {sender}\n\n{body}"
        return [Document(content, {"source": self.path, "subject": subject, "from": sender})]


# =============================================================================
# Web Loaders
# =============================================================================

class WebPageLoader(BaseLoader):
    """Load content from a web page."""
    
    def __init__(self, url: str):
        self.url = url
    
    def load(self) -> List[Document]:
        try:
            import requests
            from bs4 import BeautifulSoup
            
            response = requests.get(self.url, timeout=30)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, "html.parser")
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text(separator="\n", strip=True)
            return [Document(text, {"source": self.url})]
        except ImportError:
            raise ImportError("requests and beautifulsoup4 required")


class SitemapLoader(BaseLoader):
    """Load URLs from a sitemap."""
    
    def __init__(self, url: str, filter_urls: Optional[List[str]] = None):
        self.url = url
        self.filter_urls = filter_urls
    
    def load(self) -> List[Document]:
        try:
            import requests
            from bs4 import BeautifulSoup
            
            response = requests.get(self.url, timeout=30)
            soup = BeautifulSoup(response.text, "xml")
            
            urls = [loc.text for loc in soup.find_all("loc")]
            
            if self.filter_urls:
                urls = [u for u in urls if any(f in u for f in self.filter_urls)]
            
            docs = []
            for url in urls:
                try:
                    loader = WebPageLoader(url)
                    docs.extend(loader.load())
                except Exception:
                    continue
            return docs
        except ImportError:
            raise ImportError("requests and beautifulsoup4 required")


class YouTubeLoader(BaseLoader):
    """Load YouTube video transcripts."""
    
    def __init__(self, video_url: str):
        self.video_url = video_url
    
    def _extract_video_id(self) -> str:
        import re
        patterns = [
            r"(?:v=|\/)([\w-]{11})",
            r"youtu\.be\/([\w-]{11})",
        ]
        for pattern in patterns:
            match = re.search(pattern, self.video_url)
            if match:
                return match.group(1)
        raise ValueError(f"Could not extract video ID from {self.video_url}")
    
    def load(self) -> List[Document]:
        try:
            from youtube_transcript_api import YouTubeTranscriptApi
            
            video_id = self._extract_video_id()
            transcript = YouTubeTranscriptApi.get_transcript(video_id)
            
            content = " ".join([t["text"] for t in transcript])
            return [Document(content, {"source": self.video_url, "video_id": video_id})]
        except ImportError:
            raise ImportError("youtube-transcript-api required")


# =============================================================================
# Cloud/API Loaders
# =============================================================================

class S3Loader(BaseLoader):
    """Load files from AWS S3."""
    
    def __init__(
        self,
        bucket: str,
        key: str,
        region: str = "us-east-1",
    ):
        self.bucket = bucket
        self.key = key
        self.region = region
    
    def load(self) -> List[Document]:
        try:
            import boto3
            s3 = boto3.client("s3", region_name=self.region)
            response = s3.get_object(Bucket=self.bucket, Key=self.key)
            content = response["Body"].read().decode("utf-8")
            return [Document(content, {"source": f"s3://{self.bucket}/{self.key}"})]
        except ImportError:
            raise ImportError("boto3 required. pip install boto3")


class GCSLoader(BaseLoader):
    """Load files from Google Cloud Storage."""
    
    def __init__(self, bucket: str, blob_name: str):
        self.bucket = bucket
        self.blob_name = blob_name
    
    def load(self) -> List[Document]:
        try:
            from google.cloud import storage
            client = storage.Client()
            bucket = client.bucket(self.bucket)
            blob = bucket.blob(self.blob_name)
            content = blob.download_as_text()
            return [Document(content, {"source": f"gs://{self.bucket}/{self.blob_name}"})]
        except ImportError:
            raise ImportError("google-cloud-storage required")


class AzureBlobLoader(BaseLoader):
    """Load files from Azure Blob Storage."""
    
    def __init__(
        self,
        container: str,
        blob_name: str,
        connection_string: Optional[str] = None,
    ):
        self.container = container
        self.blob_name = blob_name
        self.connection_string = connection_string or os.getenv("AZURE_STORAGE_CONNECTION_STRING")
    
    def load(self) -> List[Document]:
        try:
            from azure.storage.blob import BlobServiceClient
            client = BlobServiceClient.from_connection_string(self.connection_string)
            container = client.get_container_client(self.container)
            blob = container.get_blob_client(self.blob_name)
            content = blob.download_blob().readall().decode("utf-8")
            return [Document(content, {"source": f"azure://{self.container}/{self.blob_name}"})]
        except ImportError:
            raise ImportError("azure-storage-blob required")


class NotionLoader(BaseLoader):
    """Load pages from Notion."""
    
    def __init__(self, page_id: str, api_key: Optional[str] = None):
        self.page_id = page_id
        self.api_key = api_key or os.getenv("NOTION_API_KEY")
    
    def load(self) -> List[Document]:
        try:
            from notion_client import Client
            notion = Client(auth=self.api_key)
            
            # Get page content
            blocks = notion.blocks.children.list(block_id=self.page_id)
            
            content_parts = []
            for block in blocks.get("results", []):
                block_type = block.get("type")
                if block_type in ["paragraph", "heading_1", "heading_2", "heading_3"]:
                    text_content = block.get(block_type, {}).get("rich_text", [])
                    for text in text_content:
                        content_parts.append(text.get("plain_text", ""))
            
            return [Document("\n".join(content_parts), {"source": f"notion://{self.page_id}"})]
        except ImportError:
            raise ImportError("notion-client required. pip install notion-client")


class GitHubLoader(BaseLoader):
    """Load files from GitHub repository."""
    
    def __init__(
        self,
        repo: str,
        path: str = "",
        branch: str = "main",
        file_filter: Optional[str] = None,
    ):
        self.repo = repo
        self.path = path
        self.branch = branch
        self.file_filter = file_filter
    
    def load(self) -> List[Document]:
        try:
            import requests
            
            api_url = f"https://api.github.com/repos/{self.repo}/contents/{self.path}"
            params = {"ref": self.branch}
            
            response = requests.get(api_url, params=params, timeout=30)
            response.raise_for_status()
            
            items = response.json()
            if not isinstance(items, list):
                items = [items]
            
            docs = []
            for item in items:
                if item["type"] == "file":
                    if self.file_filter and not item["name"].endswith(self.file_filter):
                        continue
                    
                    file_response = requests.get(item["download_url"], timeout=30)
                    content = file_response.text
                    docs.append(Document(content, {
                        "source": item["html_url"],
                        "path": item["path"],
                    }))
            
            return docs
        except ImportError:
            raise ImportError("requests required")


class ConfluenceLoader(BaseLoader):
    """Load pages from Confluence."""
    
    def __init__(
        self,
        url: str,
        space_key: str,
        username: Optional[str] = None,
        api_token: Optional[str] = None,
    ):
        self.url = url
        self.space_key = space_key
        self.username = username or os.getenv("CONFLUENCE_USERNAME")
        self.api_token = api_token or os.getenv("CONFLUENCE_API_TOKEN")
    
    def load(self) -> List[Document]:
        try:
            from atlassian import Confluence
            
            confluence = Confluence(
                url=self.url,
                username=self.username,
                password=self.api_token,
            )
            
            pages = confluence.get_all_pages_from_space(self.space_key)
            
            docs = []
            for page in pages:
                content = confluence.get_page_by_id(page["id"], expand="body.storage")
                body = content.get("body", {}).get("storage", {}).get("value", "")
                
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(body, "html.parser")
                text = soup.get_text(separator="\n", strip=True)
                
                docs.append(Document(text, {
                    "source": f"{self.url}/wiki/spaces/{self.space_key}/pages/{page['id']}",
                    "title": page.get("title", ""),
                }))
            
            return docs
        except ImportError:
            raise ImportError("atlassian-python-api required")


class SlackLoader(BaseLoader):
    """Load messages from Slack channel."""
    
    def __init__(self, channel_id: str, api_token: Optional[str] = None):
        self.channel_id = channel_id
        self.api_token = api_token or os.getenv("SLACK_API_TOKEN")
    
    def load(self) -> List[Document]:
        try:
            from slack_sdk import WebClient
            
            client = WebClient(token=self.api_token)
            response = client.conversations_history(channel=self.channel_id)
            
            messages = response.get("messages", [])
            content = "\n".join([m.get("text", "") for m in messages])
            
            return [Document(content, {"source": f"slack://{self.channel_id}"})]
        except ImportError:
            raise ImportError("slack-sdk required")


# =============================================================================
# Database Loaders
# =============================================================================

class SQLLoader(BaseLoader):
    """Load data from SQL database."""
    
    def __init__(self, query: str, connection_string: str):
        self.query = query
        self.connection_string = connection_string
    
    def load(self) -> List[Document]:
        try:
            import sqlalchemy
            engine = sqlalchemy.create_engine(self.connection_string)
            
            with engine.connect() as conn:
                result = conn.execute(sqlalchemy.text(self.query))
                rows = result.fetchall()
                columns = result.keys()
            
            docs = []
            for i, row in enumerate(rows):
                content = " | ".join([f"{col}: {val}" for col, val in zip(columns, row)])
                docs.append(Document(content, {"source": "sql", "row": i}))
            
            return docs
        except ImportError:
            raise ImportError("sqlalchemy required")


class MongoDBLoader(BaseLoader):
    """Load documents from MongoDB."""
    
    def __init__(
        self,
        connection_string: str,
        database: str,
        collection: str,
        filter: Optional[Dict] = None,
    ):
        self.connection_string = connection_string
        self.database = database
        self.collection = collection
        self.filter = filter or {}
    
    def load(self) -> List[Document]:
        try:
            from pymongo import MongoClient
            import json
            
            client = MongoClient(self.connection_string)
            db = client[self.database]
            collection = db[self.collection]
            
            docs = []
            for doc in collection.find(self.filter):
                doc.pop("_id", None)
                content = json.dumps(doc, default=str)
                docs.append(Document(content, {"source": f"mongodb://{self.database}/{self.collection}"}))
            
            return docs
        except ImportError:
            raise ImportError("pymongo required")


# =============================================================================
# Directory Loader
# =============================================================================

class DirectoryLoader(BaseLoader):
    """Load all files from a directory."""
    
    LOADER_MAPPING = {
        ".txt": TextLoader,
        ".md": MarkdownLoader,
        ".pdf": PDFLoader,
        ".docx": DOCXLoader,
        ".csv": CSVLoader,
        ".json": JSONLoader,
        ".html": HTMLLoader,
        ".py": CodeLoader,
        ".js": CodeLoader,
    }
    
    def __init__(
        self,
        path: str,
        glob: str = "**/*",
        recursive: bool = True,
        show_progress: bool = False,
    ):
        self.path = path
        self.glob = glob
        self.recursive = recursive
        self.show_progress = show_progress
    
    def load(self) -> List[Document]:
        from pathlib import Path
        
        docs = []
        base_path = Path(self.path)
        
        if self.recursive:
            files = list(base_path.rglob(self.glob))
        else:
            files = list(base_path.glob(self.glob))
        
        for file_path in files:
            if file_path.is_file():
                ext = file_path.suffix.lower()
                loader_cls = self.LOADER_MAPPING.get(ext)
                
                if loader_cls:
                    try:
                        loader = loader_cls(str(file_path))
                        docs.extend(loader.load())
                    except Exception as e:
                        if self.show_progress:
                            print(f"Error loading {file_path}: {e}")
                else:
                    # Default to text loader for unknown files
                    try:
                        loader = TextLoader(str(file_path))
                        docs.extend(loader.load())
                    except Exception:
                        pass
        
        return docs
