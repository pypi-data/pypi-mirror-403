import io
import os
import tempfile
from enum import Enum, auto

class FileType(Enum):
    """Supported file types for conversion."""
    TXT = auto()
    MD = auto()
    CSV = auto()
    DOC = auto()
    DOCX = auto()
    XLS = auto()
    XLSX = auto()
    PPT = auto()
    PPTX = auto()
    PDF = auto()
    UNKNOWN = auto()


class FileToTextConverter:
    """
    A class for converting various file types to raw text.
    """
    @staticmethod
    def mime_type_to_enum(mime_type: str) -> FileType:
        """
        Converts a MIME type to a FileType enum.

        :param mime_type: The MIME type string to convert.
        :return: The corresponding FileType enum, or UNKNOWN if not recognized.
        """
        if not mime_type or not mime_type.strip():
            return FileType.UNKNOWN

        mime_map = {
            "text/plain": FileType.TXT,
            "text/markdown": FileType.MD,
            "text/csv": FileType.CSV,
            "application/msword": FileType.DOC,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": FileType.DOCX,
            "application/vnd.ms-excel": FileType.XLS,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": FileType.XLSX,
            "application/vnd.ms-powerpoint": FileType.PPT,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": FileType.PPTX,
            "application/pdf": FileType.PDF,
        }
        return mime_map.get(mime_type, FileType.UNKNOWN)

    @staticmethod
    def file_extension_to_enum(file_path: str) -> FileType:
        """
        Converts a file path or extension to a FileType enum.

        :param file_path: The file path or extension to convert.
        :return: The corresponding FileType enum, or UNKNOWN if not recognized.
        """
        if not file_path or not file_path.strip():
            return FileType.UNKNOWN

        # Extract the file extension, removing the leading dot and making it lowercase
        file_extension = os.path.splitext(file_path)[1].lstrip('.').lower()

        ext_map = {
            "txt": FileType.TXT,
            "md": FileType.MD,
            "csv": FileType.CSV,
            "doc": FileType.DOC,
            "docx": FileType.DOCX,
            "xls": FileType.XLS,
            "xlsx": FileType.XLSX,
            "ppt": FileType.PPT,
            "pptx": FileType.PPTX,
            "pdf": FileType.PDF,
        }
        return ext_map.get(file_extension, FileType.UNKNOWN)

    @classmethod
    def parse_file(cls, file_type: FileType, file_bytes: bytes) -> str | None:
        """
        Parses file bytes based on the FileType and returns the text content.

        :param file_type: The type of the file to parse.
        :param file_bytes: The raw bytes of the file to parse.
        :return: The text content of the file, or None if the file type is not supported or parsing fails.
        """
        if file_type is None or file_bytes is None:
            return None
        if not file_bytes:
            return ""

        # Dispatch to the correct parser method
        parser_map = {
            FileType.TXT: cls._parse_plain_text,
            FileType.MD: cls._parse_plain_text,
            FileType.CSV: cls._parse_plain_text,
            FileType.DOC: cls._parse_doc,
            FileType.DOCX: cls._parse_docx,
            FileType.XLS: cls._parse_xls,
            FileType.XLSX: cls._parse_xlsx,
            FileType.PPT: cls._parse_ppt,
            FileType.PPTX: cls._parse_pptx,
            FileType.PDF: cls._parse_pdf,
        }

        parser_func = parser_map.get(file_type)

        if parser_func:
            return parser_func(file_bytes)

        return None

    @staticmethod
    def _parse_plain_text(file_bytes: bytes) -> str:
        return file_bytes.decode('utf-8')

    @staticmethod
    def _run_textract(file_bytes: bytes, extension: str) -> str:
        """
        Helper to run textract on in-memory bytes by writing to a temp file.
        Note: textract may require external system dependencies.
        """
        import textract
        with tempfile.NamedTemporaryFile(suffix=f".{extension}", delete=True) as temp_file:
            temp_file.write(file_bytes)
            temp_file.flush()  # Ensure all bytes are written to disk
            text = textract.process(temp_file.name).decode('utf-8')
        return text

    @classmethod
    def _parse_doc(cls, file_bytes: bytes) -> str:
        return cls._run_textract(file_bytes, 'doc')

    @staticmethod
    def _parse_docx(file_bytes: bytes) -> str:
        import docx
        with io.BytesIO(file_bytes) as stream:
            document = docx.Document(stream)
            return "\n".join(para.text for para in document.paragraphs if para.text.strip())

    @staticmethod
    def _parse_xls(file_bytes: bytes) -> str:
        import xlrd
        workbook = xlrd.open_workbook(file_contents=file_bytes)
        text_parts = []
        for sheet in workbook.sheets():
            text_parts.append(f"Sheet: {sheet.name}\n")
            for row_idx in range(sheet.nrows):
                row_cells = []
                for col_idx in range(sheet.ncols):
                    cell_text = str(sheet.cell_value(row_idx, col_idx))
                    if cell_text.strip():
                        row_cells.append(cell_text + "\t")
                if row_cells:
                    text_parts.append("".join(row_cells))
                text_parts.append("\n")
            text_parts.append("\n")
        return "".join(text_parts)

    @staticmethod
    def _parse_xlsx(file_bytes: bytes) -> str:
        import openpyxl
        with io.BytesIO(file_bytes) as stream:
            workbook = openpyxl.load_workbook(stream, read_only=True)
            text_parts = []
            for sheet in workbook.worksheets:
                text_parts.append(f"Sheet: {sheet.title}\n")
                for row in sheet.iter_rows():
                    row_cells = []
                    for cell in row:
                        cell_text = str(cell.value) if cell.value is not None else ""
                        if cell_text.strip():
                            row_cells.append(cell_text + "\t")
                    if row_cells:
                        text_parts.append("".join(row_cells))
                    text_parts.append("\n")
                text_parts.append("\n")
            return "".join(text_parts)

    @classmethod
    def _parse_ppt(cls, file_bytes: bytes) -> str:
        return cls._run_textract(file_bytes, 'ppt')

    @staticmethod
    def _parse_pptx(file_bytes: bytes) -> str:
        import pptx
        with io.BytesIO(file_bytes) as stream:
            presentation = pptx.Presentation(stream)
            text_parts = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        if text and text.strip():
                            text_parts.append(text)
            return "\n".join(text_parts)

    @staticmethod
    def _parse_pdf(file_bytes: bytes) -> str:
        """Parses a PDF file's bytes and extracts text using PyMuPDF."""
        import pymupdf
        text_parts = []
        with io.BytesIO(file_bytes) as stream:
            with pymupdf.open(stream=stream) as doc:
                for page in doc:
                    text_parts.append(page.get_text())
        return "\n".join(text_parts)
