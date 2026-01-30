"""Chart layout renderer."""

from typing import Any, Dict, List

from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from slidegen.layouts.base import BaseLayoutRenderer


class ChartLayoutRenderer(BaseLayoutRenderer):
    """Renderer for chart slide layout."""

    # Chart type mapping
    CHART_TYPE_MAP = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }

    def render(
        self, slide: Slide, data: Dict[str, Any], presentation: Presentation
    ) -> None:
        """
        Render a chart slide.
        
        Args:
            slide: The slide to render to
            data: Slide data with 'title' and 'chart' fields
            presentation: The parent presentation
        """
        title = data.get("title", "")
        chart_data = data.get("chart", {})

        # Title positioning (top area)
        slide_width = self.SLIDE_WIDTH
        title_width = slide_width - (2 * self.MARGIN)
        title_left = self.MARGIN
        title_top = self.TITLE_TOP

        # Use chart title if provided, otherwise slide title
        chart_title = chart_data.get("title") or title

        if chart_title:
            self.add_text_box(
                slide=slide,
                left=title_left,
                top=title_top,
                width=title_width,
                height=self.TITLE_HEIGHT,
                text=chart_title,
                font_size=32,
                bold=True,
                align=PP_ALIGN.LEFT,
            )

        # Chart positioning (content area)
        chart_top = self.CONTENT_TOP if chart_title else self.TITLE_TOP
        chart_width = slide_width - (2 * self.MARGIN)
        chart_left = self.MARGIN
        chart_height = self.CONTENT_HEIGHT

        # Get chart type
        chart_type_str = chart_data.get("type", "column")
        chart_type = self.CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Get chart data
        data_spec = chart_data.get("data", {})
        
        # Handle inline data
        if "labels" in data_spec and "values" in data_spec:
            labels = data_spec["labels"]
            values = data_spec["values"]
            series = data_spec.get("series", [])
            
            # Add chart to slide
            chart_shape = slide.shapes.add_chart(
                chart_type,
                Inches(chart_left),
                Inches(chart_top),
                Inches(chart_width),
                Inches(chart_height),
            )
            
            chart = chart_shape.chart
            chart.has_legend = len(series) > 0
            
            # Set chart data
            if series:
                # Multi-series chart
                self._set_multi_series_chart_data(chart, labels, series)
            else:
                # Single series chart
                self._set_single_series_chart_data(chart, labels, values)
            
            # Set axis labels
            x_axis_label = chart_data.get("x_axis_label")
            y_axis_label = chart_data.get("y_axis_label")
            if x_axis_label and hasattr(chart, "category_axis"):
                chart.category_axis.axis_title.text_frame.text = x_axis_label
            if y_axis_label and hasattr(chart, "value_axis"):
                chart.value_axis.axis_title.text_frame.text = y_axis_label
        else:
            # Handle external data source (CSV/JSON)
            # For now, we'll skip this and log a warning
            # Full implementation would require data loading module
            pass

    def _set_single_series_chart_data(
        self, chart, labels: List[str], values: List[float]
    ) -> None:
        """
        Set data for a single-series chart.
        
        Args:
            chart: The chart object
            labels: Category labels
            values: Data values
        """
        chart_data = chart.plots[0]
        chart_data.categories = labels
        chart_data.series[0].values = values

    def _set_multi_series_chart_data(
        self, chart, labels: List[str], series: List[Dict[str, Any]]
    ) -> None:
        """
        Set data for a multi-series chart.
        
        Args:
            chart: The chart object
            labels: Category labels
            series: List of series data (each with 'name' and 'values')
        """
        chart_data = chart.plots[0]
        chart_data.categories = labels
        
        # Clear existing series and add new ones
        while len(chart_data.series) > 0:
            chart_data.series[0].delete()
        
        for series_item in series:
            series_name = series_item.get("name", "")
            series_values = series_item.get("values", [])
            new_series = chart_data.series.add(series_name)
            new_series.values = series_values

