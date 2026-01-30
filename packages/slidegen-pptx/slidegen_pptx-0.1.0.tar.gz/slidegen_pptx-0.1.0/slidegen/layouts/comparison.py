"""Comparison layout renderer."""

from typing import Any, Dict

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from slidegen.layouts.base import BaseLayoutRenderer


class ComparisonLayoutRenderer(BaseLayoutRenderer):
    """Renderer for comparison slide layout (before/after, option A/B)."""

    COLUMN_GAP = Inches(0.3)  # Gap between columns

    def render(
        self, slide: Slide, data: Dict[str, Any], presentation: Presentation
    ) -> None:
        """
        Render a comparison slide.
        
        Args:
            slide: The slide to render to
            data: Slide data with 'title', 'before'/'after' or 'left'/'right' fields
            presentation: The parent presentation
        """
        title = data.get("title", "")
        before_content = data.get("before", {})
        after_content = data.get("after", {})
        
        # Fallback to left/right if before/after not provided
        if not before_content:
            before_content = data.get("left", {})
        if not after_content:
            after_content = data.get("right", {})

        # Title positioning (top area)
        slide_width = self.SLIDE_WIDTH
        title_width = slide_width - (2 * self.MARGIN)
        title_left = self.MARGIN
        title_top = self.TITLE_TOP

        if title:
            self.add_text_box(
                slide=slide,
                left=title_left,
                top=title_top,
                width=title_width,
                height=self.TITLE_HEIGHT,
                text=title,
                font_size=32,
                bold=True,
                align=PP_ALIGN.LEFT,
            )

        # Column positioning
        content_top = self.CONTENT_TOP if title else self.TITLE_TOP
        content_height = self.CONTENT_HEIGHT
        available_width = slide_width - (2 * self.MARGIN) - self.COLUMN_GAP
        column_width = available_width / 2
        column_left = self.MARGIN
        column_right = column_left + column_width + self.COLUMN_GAP

        # Add column labels
        label_height = Inches(0.4)
        before_label = before_content.get("label", "Before")
        after_label = after_content.get("label", "After")
        
        # Before column label
        if before_content:
            self.add_text_box(
                slide=slide,
                left=column_left,
                top=content_top,
                width=column_width,
                height=label_height,
                text=before_label,
                font_size=20,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
        
        # After column label
        if after_content:
            self.add_text_box(
                slide=slide,
                left=column_right,
                top=content_top,
                width=column_width,
                height=label_height,
                text=after_label,
                font_size=20,
                bold=True,
                align=PP_ALIGN.CENTER,
            )

        # Render before column content
        if before_content:
            content_data = before_content.get("content", before_content)
            self._render_column_content(
                slide=slide,
                left=column_left,
                top=content_top + label_height,
                width=column_width,
                height=content_height - label_height,
                content=content_data,
            )

        # Render after column content
        if after_content:
            content_data = after_content.get("content", after_content)
            self._render_column_content(
                slide=slide,
                left=column_right,
                top=content_top + label_height,
                width=column_width,
                height=content_height - label_height,
                content=content_data,
            )

    def _render_column_content(
        self,
        slide: Slide,
        left: float,
        top: float,
        width: float,
        height: float,
        content: Dict[str, Any],
    ) -> None:
        """
        Render content in a column.
        
        Args:
            slide: The slide to render to
            left: Left position in inches
            top: Top position in inches
            width: Column width in inches
            height: Column height in inches
            content: Content data (can be bullet_list, text, etc.)
        """
        content_type = content.get("type", "text")
        
        if content_type == "bullet_list":
            bullets = content.get("bullets", [])
            if bullets:
                self._render_bullets_in_column(
                    slide=slide,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    bullets=bullets,
                )
        elif content_type == "text":
            text = content.get("text", "")
            if text:
                self.add_text_box(
                    slide=slide,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    text=text,
                    font_size=16,
                    bold=False,
                    align=PP_ALIGN.LEFT,
                )
        else:
            # Default to text if type is unknown
            text = str(content.get("text", ""))
            if text:
                self.add_text_box(
                    slide=slide,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    text=text,
                    font_size=16,
                    bold=False,
                    align=PP_ALIGN.LEFT,
                )

    def _render_bullets_in_column(
        self,
        slide: Slide,
        left: float,
        top: float,
        width: float,
        height: float,
        bullets: list,
    ) -> None:
        """
        Render bullet list in a column.
        
        Args:
            slide: The slide to render to
            left: Left position in inches
            top: Top position in inches
            width: Column width in inches
            height: Column height in inches
            bullets: List of bullet items
        """
        max_bullets = len(bullets)
        if max_bullets == 0:
            return

        bullet_height = min(
            height / max_bullets if max_bullets > 0 else height,
            Inches(0.5)  # Minimum spacing
        )

        for i, bullet_item in enumerate(bullets):
            # Handle both string and object format
            if isinstance(bullet_item, str):
                bullet_text = bullet_item
                bullet_level = 0
            else:
                bullet_text = bullet_item.get("text", "")
                bullet_level = bullet_item.get("level", 0)

            if not bullet_text:
                continue

            # Calculate position
            bullet_top = top + (i * bullet_height)
            indent = Inches(0.3) + (bullet_level * Inches(0.3))
            bullet_left = left + indent
            bullet_width = width - indent

            # Add bullet point
            self._add_bullet_point(
                slide=slide,
                left=bullet_left,
                top=bullet_top,
                width=bullet_width,
                height=bullet_height,
                text=bullet_text,
                level=bullet_level,
            )

    def _add_bullet_point(
        self,
        slide: Slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        level: int = 0,
    ) -> None:
        """
        Add a bullet point to the slide.
        
        Args:
            slide: The slide to add the bullet to
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            text: Bullet text
            level: Bullet nesting level (0 or 1)
        """
        font_size = 16 if level == 0 else 14

        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.25)  # Space for bullet
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)

        # Set paragraph formatting
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = False
        paragraph.level = level

        # Add bullet character
        paragraph.bullet = True

