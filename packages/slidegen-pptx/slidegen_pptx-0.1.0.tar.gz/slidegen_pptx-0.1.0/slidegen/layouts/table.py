"""Table layout renderer."""

from typing import Any, Dict, List

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from slidegen.layouts.base import BaseLayoutRenderer


class TableLayoutRenderer(BaseLayoutRenderer):
    """Renderer for table slide layout."""

    def render(
        self, slide: Slide, data: Dict[str, Any], presentation: Presentation
    ) -> None:
        """
        Render a table slide.
        
        Args:
            slide: The slide to render to
            data: Slide data with 'title' and 'table' fields
            presentation: The parent presentation
        """
        title = data.get("title", "")
        table_data = data.get("table", {})

        # Title positioning (top area)
        slide_width = self.SLIDE_WIDTH
        title_width = slide_width - (2 * self.MARGIN)
        title_left = self.MARGIN
        title_top = self.TITLE_TOP

        if title:
            self.add_text_box(
                slide=slide,
                left=title_left,
                top=title_top,
                width=title_width,
                height=self.TITLE_HEIGHT,
                text=title,
                font_size=32,
                bold=True,
                align=PP_ALIGN.LEFT,
            )

        # Table positioning (content area)
        table_top = self.CONTENT_TOP if title else self.TITLE_TOP
        table_width = slide_width - (2 * self.MARGIN)
        table_left = self.MARGIN
        table_height = self.CONTENT_HEIGHT

        # Get table data
        data_spec = table_data.get("data", [])
        header_row = table_data.get("header_row", True)

        # Handle inline data (array of rows)
        if isinstance(data_spec, list) and len(data_spec) > 0:
            rows = data_spec
            num_rows = len(rows)
            num_cols = len(rows[0]) if rows else 0

            if num_rows > 0 and num_cols > 0:
                # Add table to slide
                table_shape = slide.shapes.add_table(
                    num_rows,
                    num_cols,
                    Inches(table_left),
                    Inches(table_top),
                    Inches(table_width),
                    Inches(table_height),
                )
                table = table_shape.table

                # Populate table
                for row_idx, row_data in enumerate(rows):
                    for col_idx, cell_value in enumerate(row_data):
                        if col_idx < num_cols:
                            cell = table.cell(row_idx, col_idx)
                            cell.text = str(cell_value)

                            # Style header row if specified
                            if header_row and row_idx == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = None  # Use default header color
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.bold = True
                                    paragraph.font.size = Pt(14)
                            else:
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(12)

                # Auto-fit column widths
                self._auto_fit_columns(table, num_cols, table_width)
        else:
            # Handle external data source (CSV)
            # For now, we'll skip this and log a warning
            # Full implementation would require data loading module
            pass

    def _auto_fit_columns(self, table, num_cols: int, total_width: float) -> None:
        """
        Auto-fit column widths evenly.
        
        Args:
            table: The table object
            num_cols: Number of columns
            total_width: Total table width in inches
        """
        if num_cols == 0:
            return

        col_width = total_width / num_cols
        for col_idx in range(num_cols):
            for cell in table.columns[col_idx].cells:
                cell.width = Inches(col_width)

