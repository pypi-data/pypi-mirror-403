"""Image layout renderer."""

import os
from pathlib import Path
from typing import Any, Dict

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.presentation import Presentation
from pptx.util import Inches

from slidegen.layouts.base import BaseLayoutRenderer


class ImageLayoutRenderer(BaseLayoutRenderer):
    """Renderer for image slide layout."""

    def render(
        self, slide: Slide, data: Dict[str, Any], presentation: Presentation
    ) -> None:
        """
        Render an image slide.
        
        Args:
            slide: The slide to render to
            data: Slide data with 'title' and 'image' fields
            presentation: The parent presentation
        """
        title = data.get("title", "")
        image_data = data.get("image", {})

        # Title positioning (top area)
        slide_width = self.SLIDE_WIDTH
        title_width = slide_width - (2 * self.MARGIN)
        title_left = self.MARGIN
        title_top = self.TITLE_TOP

        if title:
            self.add_text_box(
                slide=slide,
                left=title_left,
                top=title_top,
                width=title_width,
                height=self.TITLE_HEIGHT,
                text=title,
                font_size=32,
                bold=True,
                align=PP_ALIGN.LEFT,
            )

        # Image positioning (content area)
        image_top = self.CONTENT_TOP if title else self.TITLE_TOP
        image_width = slide_width - (2 * self.MARGIN)
        image_left = self.MARGIN
        image_height = self.CONTENT_HEIGHT

        # Get image path
        image_src = image_data.get("src", "")
        if not image_src:
            return

        # Resolve image path (relative to schema file or absolute)
        # For now, we'll try to load from the current directory or absolute path
        # Full implementation would need schema file path context
        image_path = Path(image_src)
        if not image_path.is_absolute():
            # Try current directory
            image_path = Path.cwd() / image_src

        if not image_path.exists():
            # Image not found, skip
            return

        # Calculate image dimensions to fit in available space while maintaining aspect ratio
        try:
            from PIL import Image as PILImage
            
            with PILImage.open(image_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
                
                # Fit to available space
                available_width = image_width
                available_height = image_height
                
                # Calculate dimensions that fit
                if available_width / available_height > aspect_ratio:
                    # Height is limiting factor
                    display_height = available_height
                    display_width = display_height * aspect_ratio
                else:
                    # Width is limiting factor
                    display_width = available_width
                    display_height = display_width / aspect_ratio
                
                # Center the image
                image_left_offset = image_left + (image_width - display_width) / 2
                image_top_offset = image_top + (image_height - display_height) / 2
                
                # Add image to slide
                slide.shapes.add_picture(
                    str(image_path),
                    Inches(image_left_offset),
                    Inches(image_top_offset),
                    Inches(display_width),
                    Inches(display_height),
                )
        except ImportError:
            # PIL not available, use default size
            slide.shapes.add_picture(
                str(image_path),
                Inches(image_left),
                Inches(image_top),
                Inches(image_width),
                Inches(image_height),
            )
        except Exception:
            # Image loading failed, skip
            pass

        # Add caption if provided
        caption = image_data.get("caption", "")
        if caption:
            caption_top = image_top + image_height + Inches(0.1)
            caption_height = Inches(0.3)
            self.add_text_box(
                slide=slide,
                left=image_left,
                top=caption_top,
                width=image_width,
                height=caption_height,
                text=caption,
                font_size=12,
                bold=False,
                align=PP_ALIGN.CENTER,
            )

