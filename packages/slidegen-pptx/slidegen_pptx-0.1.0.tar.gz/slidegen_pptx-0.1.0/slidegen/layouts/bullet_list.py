"""Bullet list layout renderer."""

from typing import Any, Dict, List, Union

from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.slide import Slide
from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from slidegen.layouts.base import BaseLayoutRenderer


class BulletListLayoutRenderer(BaseLayoutRenderer):
    """Renderer for bullet list slide layout."""

    def render(
        self, slide: Slide, data: Dict[str, Any], presentation: Presentation
    ) -> None:
        """
        Render a bullet list slide.
        
        Args:
            slide: The slide to render to
            data: Slide data with 'title' and 'bullets' array
            presentation: The parent presentation
        """
        title = data.get("title", "")
        bullets = data.get("bullets", [])

        # Title positioning (top area)
        if title:
            self.add_text_box(
                slide=slide,
                left=self.MARGIN,
                top=self.TITLE_TOP,
                width=self.SLIDE_WIDTH - (2 * self.MARGIN),
                height=self.TITLE_HEIGHT,
                text=title,
                font_size=32,
                bold=True,
                align=PP_ALIGN.LEFT,
            )

        # Bullet list positioning (content area)
        if bullets:
            content_top = self.CONTENT_TOP if title else self.TITLE_TOP
            content_height = self.CONTENT_HEIGHT if title else (
                self.SLIDE_HEIGHT - self.TITLE_TOP - self.MARGIN
            )
            
            self._render_bullets(
                slide=slide,
                bullets=bullets,
                left=self.MARGIN,
                top=content_top,
                width=self.SLIDE_WIDTH - (2 * self.MARGIN),
                height=content_height,
            )

    def _render_bullets(
        self,
        slide: Slide,
        bullets: List[Union[str, Dict[str, Any]]],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        """
        Render a list of bullets with automatic text flow and overflow handling.
        
        Args:
            slide: The slide to render to
            bullets: List of bullet items (strings or dicts with text/level)
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
        """
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)

        # Clear default paragraph
        text_frame.clear()
        
        # Add bullets
        for bullet_item in bullets:
            if isinstance(bullet_item, str):
                bullet_text = bullet_item
                level = 0
            else:
                bullet_text = bullet_item.get("text", "")
                level = bullet_item.get("level", 0)
            
            if not bullet_text:
                continue
            
            # Add paragraph for this bullet
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet_text
            paragraph.level = level
            paragraph.font.size = Pt(18 if level == 0 else 16)
            paragraph.space_after = Pt(6)
            
            # Set indentation based on level
            if level == 1:
                paragraph.left_indent = Inches(0.5)
                paragraph.first_line_indent = Inches(-0.25)

