"""Core SlideGenerator class for rendering presentations."""

import json
from pathlib import Path
from typing import Any, Dict, Optional

import yaml
from pptx import Presentation
from pptx.util import Inches

from slidegen.core.registry import LayoutRegistry
from slidegen.layouts import (
    BlankLayoutRenderer,
    BulletListLayoutRenderer,
    ChartLayoutRenderer,
    ComparisonLayoutRenderer,
    ImageLayoutRenderer,
    QuoteLayoutRenderer,
    SectionHeaderLayoutRenderer,
    TableLayoutRenderer,
    TitleLayoutRenderer,
    TwoColumnLayoutRenderer,
)
from slidegen.theming import Theme, load_theme
from slidegen.validator import validate


def _register_default_layouts(registry: LayoutRegistry) -> None:
    """Register default layout renderers."""
    registry.register("title", TitleLayoutRenderer())
    registry.register("section_header", SectionHeaderLayoutRenderer())
    registry.register("bullet_list", BulletListLayoutRenderer())
    registry.register("two_column", TwoColumnLayoutRenderer())
    registry.register("chart", ChartLayoutRenderer())
    registry.register("table", TableLayoutRenderer())
    registry.register("comparison", ComparisonLayoutRenderer())
    registry.register("image", ImageLayoutRenderer())
    registry.register("quote", QuoteLayoutRenderer())
    registry.register("blank", BlankLayoutRenderer())


class SlideGenerator:
    """
    Main class for generating PowerPoint presentations from schemas.
    
    Usage:
        gen = SlideGenerator()
        gen.from_schema("deck.yaml")
        gen.build("output.pptx")
    """

    def __init__(self, theme: Optional[str] = None, theme_name: Optional[str] = None):
        """
        Initialize a SlideGenerator.
        
        Args:
            theme: Optional path to PowerPoint template file
            theme_name: Optional theme name (default, corporate, dark)
        """
        self.theme_path = theme
        self.theme_name = theme_name
        self.theme_obj: Optional[Theme] = None
        if theme or theme_name:
            self.theme_obj = load_theme(theme_path=theme, theme_name=theme_name)
        else:
            self.theme_obj = Theme.default()
        self.schema: Optional[Dict[str, Any]] = None
        self.presentation: Optional[Presentation] = None
        self.registry = LayoutRegistry()
        _register_default_layouts(self.registry)

    def from_schema(self, schema_path: str) -> None:
        """
        Load a schema from a YAML or JSON file.
        
        Args:
            schema_path: Path to schema file
            
        Raises:
            FileNotFoundError: If schema file doesn't exist
            ValueError: If schema is invalid
        """
        path = Path(schema_path)
        if not path.exists():
            raise FileNotFoundError(f"Schema file not found: {schema_path}")

        # Load schema file
        with open(path, "r", encoding="utf-8") as f:
            if path.suffix in [".yaml", ".yml"]:
                self.schema = yaml.safe_load(f)
            elif path.suffix == ".json":
                self.schema = json.load(f)
            else:
                raise ValueError(f"Unsupported file format: {path.suffix}")

        # Validate schema
        result = validate(self.schema)
        if not result.is_valid:
            error_messages = "\n".join(
                f"- {err.get('field', 'unknown')}: {err.get('message', 'unknown error')}"
                for err in result.errors
            )
            raise ValueError(f"Schema validation failed:\n{error_messages}")

    def from_dict(self, schema: Dict[str, Any]) -> None:
        """
        Load a schema from a dictionary.
        
        Args:
            schema: Schema dictionary
            
        Raises:
            ValueError: If schema is invalid
        """
        # Validate schema
        result = validate(schema)
        if not result.is_valid:
            error_messages = "\n".join(
                f"- {err.get('field', 'unknown')}: {err.get('message', 'unknown error')}"
                for err in result.errors
            )
            raise ValueError(f"Schema validation failed:\n{error_messages}")

        self.schema = schema

    def build(self, output_path: str) -> None:
        """
        Render the loaded schema to a PowerPoint file.
        
        Args:
            output_path: Path where the .pptx file should be saved
            
        Raises:
            ValueError: If no schema has been loaded
        """
        if self.schema is None:
            raise ValueError("No schema loaded. Call from_schema() or from_dict() first.")

        # Initialize presentation
        if self.theme_path and Path(self.theme_path).exists():
            self.presentation = Presentation(self.theme_path)
        else:
            self.presentation = Presentation()
            # Set default slide size (16:9 widescreen)
            self.presentation.slide_width = Inches(10)
            self.presentation.slide_height = Inches(5.625)
            
            # Apply theme colors if theme is set
            if self.theme_obj:
                # Set background color for all slides (via slide master)
                if self.theme_obj.background_color:
                    bg_color = self.theme_obj.get_background_rgb()
                    if bg_color:
                        # Apply to slide master background
                        slide_master = self.presentation.slide_masters[0]
                        slide_master.background.fill.solid()
                        slide_master.background.fill.fore_color.rgb = bg_color

        # Get presentation data
        presentation_data = self.schema.get("presentation", {})
        slides_data = presentation_data.get("slides", [])

        # Render each slide
        for slide_data in slides_data:
            layout_type = slide_data.get("layout")
            if not layout_type:
                continue

            # Get layout renderer
            renderer = self.registry.get_renderer(layout_type)
            if renderer is None:
                # Skip unknown layouts (will be handled by validation)
                continue

            # Create new slide
            slide = self.presentation.slides.add_slide(
                self.presentation.slide_layouts[6]  # Blank layout
            )

            # Render the slide
            renderer.render(slide, slide_data, self.presentation)

        # Save presentation
        output = Path(output_path)
        self.presentation.save(str(output))
