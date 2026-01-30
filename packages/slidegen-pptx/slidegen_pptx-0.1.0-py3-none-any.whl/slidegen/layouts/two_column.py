"""Two column layout renderer."""

from typing import Any, Dict

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from slidegen.layouts.base import BaseLayoutRenderer


class TwoColumnLayoutRenderer(BaseLayoutRenderer):
    """Renderer for two column slide layout."""

    COLUMN_GAP = Inches(0.3)  # Gap between columns
    COLUMN_WIDTH = (BaseLayoutRenderer.SLIDE_WIDTH - (2 * BaseLayoutRenderer.MARGIN) - Inches(0.3)) / 2

    def render(
        self, slide: Slide, data: Dict[str, Any], presentation: Presentation
    ) -> None:
        """
        Render a two column slide.
        
        Args:
            slide: The slide to render to
            data: Slide data with 'title', 'left', and 'right' columns
            presentation: The parent presentation
        """
        title = data.get("title", "")
        left_content = data.get("left", {})
        right_content = data.get("right", {})

        # Title positioning (top area)
        if title:
            self.add_text_box(
                slide=slide,
                left=self.MARGIN,
                top=self.TITLE_TOP,
                width=self.SLIDE_WIDTH - (2 * self.MARGIN),
                height=self.TITLE_HEIGHT,
                text=title,
                font_size=32,
                bold=True,
                align=PP_ALIGN.LEFT,
            )

        # Column positioning
        content_top = self.CONTENT_TOP if title else self.TITLE_TOP
        content_height = self.CONTENT_HEIGHT if title else (
            self.SLIDE_HEIGHT - self.TITLE_TOP - self.MARGIN
        )
        
        left_x = self.MARGIN
        right_x = self.MARGIN + self.COLUMN_WIDTH + self.COLUMN_GAP

        # Render left column
        if left_content:
            self._render_column(
                slide=slide,
                content=left_content,
                left=left_x,
                top=content_top,
                width=self.COLUMN_WIDTH,
                height=content_height,
            )

        # Render right column
        if right_content:
            self._render_column(
                slide=slide,
                content=right_content,
                left=right_x,
                top=content_top,
                width=self.COLUMN_WIDTH,
                height=content_height,
            )

    def _render_column(
        self,
        slide: Slide,
        content: Dict[str, Any],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        """
        Render content in a column.
        
        Args:
            slide: The slide to render to
            content: Column content dict with 'type' and content data
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
        """
        content_type = content.get("type", "text")
        
        if content_type == "bullet_list":
            bullets = content.get("bullets", [])
            self._render_bullets_in_column(slide, bullets, left, top, width, height)
        elif content_type == "text":
            text = content.get("text", "")
            self._render_text_in_column(slide, text, left, top, width, height)
        elif content_type == "image":
            # Image rendering will be handled in TASK-008
            # For now, just add a placeholder text
            self.add_text_box(
                slide=slide,
                left=left,
                top=top,
                width=width,
                height=height,
                text=f"[Image: {content.get('image', {}).get('src', 'placeholder')}]",
                font_size=14,
                bold=False,
                align=PP_ALIGN.CENTER,
            )

    def _render_bullets_in_column(
        self,
        slide: Slide,
        bullets: list,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        """Render bullets in a column."""
        from slidegen.layouts.bullet_list import BulletListLayoutRenderer
        
        bullet_renderer = BulletListLayoutRenderer()
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)

        # Clear default paragraph
        text_frame.clear()
        
        # Add bullets
        for bullet_item in bullets:
            if isinstance(bullet_item, str):
                bullet_text = bullet_item
                level = 0
            else:
                bullet_text = bullet_item.get("text", "")
                level = bullet_item.get("level", 0)
            
            if not bullet_text:
                continue
            
            # Add paragraph for this bullet
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet_text
            paragraph.level = level
            paragraph.font.size = Pt(16 if level == 0 else 14)
            paragraph.space_after = Pt(6)
            
            # Set indentation based on level
            if level == 1:
                paragraph.left_indent = Inches(0.4)
                paragraph.first_line_indent = Inches(-0.2)

    def _render_text_in_column(
        self,
        slide: Slide,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        """Render plain text in a column."""
        self.add_text_box(
            slide=slide,
            left=left,
            top=top,
            width=width,
            height=height,
            text=text,
            font_size=16,
            bold=False,
            align=PP_ALIGN.LEFT,
        )

