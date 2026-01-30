"""Base classes for layout renderers."""

from typing import Any, Dict

from pptx.slide import Slide
from pptx.presentation import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class BaseLayoutRenderer:
    """Base class for layout renderers."""

    # Standard positioning grid (in inches)
    SLIDE_WIDTH = Inches(10)
    SLIDE_HEIGHT = Inches(5.625)
    MARGIN = Inches(0.5)
    TITLE_TOP = Inches(0.5)
    TITLE_HEIGHT = Inches(0.8)
    CONTENT_TOP = Inches(1.5)
    CONTENT_HEIGHT = Inches(3.5)

    def render(self, slide: Slide, data: Dict[str, Any], presentation: Presentation) -> None:
        """
        Render a slide with the given data.
        
        Args:
            slide: The pptx Slide object to render to
            data: Slide data from schema
            presentation: The parent Presentation object
        """
        raise NotImplementedError("Subclasses must implement render()")

    def add_text_box(
        self,
        slide: Slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_size: int = 18,
        bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ) -> None:
        """
        Add a text box to a slide.
        
        Args:
            slide: The slide to add the text box to
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            text: Text content
            font_size: Font size in points
            bold: Whether text should be bold
            align: Text alignment
        """
        from pptx.shapes.autoshape import Shape
        from pptx.enum.shapes import MSO_SHAPE

        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True

        # Set paragraph formatting
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold

