"""Unified document interrogation: retrieve, screenshot, search."""
__all__ = ['open_document', 'DocumentHandle', 'DocRouterError', 'UnsupportedFileTypeError',
           'UnsupportedOperationError', 'DocumentNotOpenError', 'ParseError']
__version__ = '0.1.0'

import io
import re
import tempfile
import uuid
from pathlib import Path
from typing import BinaryIO, Union

# === Exceptions ===
class DocRouterError(Exception): "Base for all docrouter errors."
class UnsupportedFileTypeError(DocRouterError): "File format not supported."
class UnsupportedOperationError(DocRouterError): "Operation not supported for this format."
class DocumentNotOpenError(DocRouterError): "Document ID not found in registry."
class ParseError(DocRouterError): "Failed to parse document."

# === MIME types ===
_MIMES = {
    'pdf': 'application/pdf', 'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'html': 'text/html', 'markdown': 'text/markdown', 'text': 'text/plain',
    'image': 'image/png', 'unknown': 'application/octet-stream'
}

# === Utilities ===
def _chunk(text: str, target: int = 2000) -> list[str]:
    """Split text into ~target char chunks, aligned to paragraph boundaries."""
    if not text:
        return ['']

    paragraphs = re.split(r'\n\s*\n', text)
    chunks = []
    current_chunk = []
    current_size = 0

    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if not paragraph:
            continue
        paragraph_len = len(paragraph) + 2  # +2 for '\n\n' separator

        if current_size + paragraph_len > target and current_chunk:
            chunks.append('\n\n'.join(current_chunk))
            current_chunk = []
            current_size = 0

        current_chunk.append(paragraph)
        current_size += paragraph_len

    if current_chunk:
        chunks.append('\n\n'.join(current_chunk))

    return chunks if chunks else ['']

def _resolve_input(inp, filename: str | None) -> tuple[Path, bytes]:
    """Resolve input to (path, raw_bytes). Creates temp file if needed."""
    if isinstance(inp, (str, Path)):
        p = Path(inp)
        return p, p.read_bytes()
    elif isinstance(inp, bytes):
        if not filename: raise ValueError("filename required when input is bytes")
        tmp = Path(tempfile.gettempdir()) / f"docrouter_{uuid.uuid4().hex[:8]}_{filename}"
        tmp.write_bytes(inp)
        return tmp, inp
    elif hasattr(inp, 'read'):
        raw = inp.read()
        if not filename: raise ValueError("filename required for file-like input")
        tmp = Path(tempfile.gettempdir()) / f"docrouter_{uuid.uuid4().hex[:8]}_{filename}"
        tmp.write_bytes(raw)
        return tmp, raw
    else:
        raise TypeError(f"input must be str, Path, bytes, or file-like, got {type(inp)}")

# === Base Doc class ===
class Doc:
    """Base for document parsers. Subclasses override as needed."""
    fmt: str = 'unknown'
    unit_type: str = 'chunk'

    def __init__(self, raw: bytes, path: Path):
        self._raw = raw
        self._path = path

    def text(self) -> str:
        return ''

    def units(self) -> list[str]:
        return _chunk(self.text())

    def meta(self) -> dict:
        return {}

    def render(self, idx: int, dpi: int) -> bytes:
        raise UnsupportedOperationError(f"render_page only supported for PDF, not {self.fmt}")

# === TxtDoc ===
class TxtDoc(Doc):
    fmt = 'text'
    unit_type = 'chunk'

    def __init__(self, raw: bytes, path: Path):
        super().__init__(raw, path)
        self._text = raw.decode('utf-8', errors='replace')

    def text(self) -> str:
        return self._text

# === MdDoc ===
class MdDoc(Doc):
    fmt = 'markdown'
    unit_type = 'chunk'

    def __init__(self, raw: bytes, path: Path):
        super().__init__(raw, path)
        self._text = raw.decode('utf-8', errors='replace')

    def text(self) -> str:
        return self._text

    def meta(self) -> dict:
        for line in self._text.split('\n'):
            if line.startswith('# '):
                return {'title': line[2:].strip()}
        return {}

# === PDFDoc ===
class PDFDoc(Doc):
    fmt = 'pdf'
    unit_type = 'page'

    def __init__(self, raw: bytes, path: Path):
        super().__init__(raw, path)
        try:
            import pymupdf
        except ImportError:
            raise ImportError("pymupdf required for PDF support: pip install docrouter[pdf]")
        self._doc = pymupdf.open(stream=raw, filetype='pdf')

    def text(self) -> str:
        return '\n\n'.join(self.units())

    def units(self) -> list[str]:
        return [page.get_text() for page in self._doc]

    def meta(self) -> dict:
        metadata = self._doc.metadata or {}
        result = {}
        if metadata.get('title'):
            result['title'] = metadata['title']
        if metadata.get('author'):
            result['author'] = metadata['author']
        return result

    def render(self, idx: int, dpi: int) -> bytes:
        if not 0 <= idx < len(self._doc):
            raise IndexError(f"page {idx} out of range (0-{len(self._doc)-1})")
        pix = self._doc[idx].get_pixmap(dpi=dpi)
        return pix.tobytes('png')

# === DocxDoc ===
class DocxDoc(Doc):
    fmt = 'docx'
    unit_type = 'section'

    def __init__(self, raw: bytes, path: Path):
        super().__init__(raw, path)
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx required for DOCX support: pip install docrouter[docx]")
        self._doc = Document(io.BytesIO(raw))

    def text(self) -> str:
        parts = []
        for paragraph in self._doc.paragraphs:
            if paragraph.text.strip():
                parts.append(paragraph.text)
        for table in self._doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(' | '.join(cells))
        return '\n\n'.join(parts)

    def meta(self) -> dict:
        props = self._doc.core_properties
        result = {}
        if props.title:
            result['title'] = props.title
        if props.author:
            result['author'] = props.author
        return result

# === PptxDoc ===
class PptxDoc(Doc):
    fmt = 'pptx'
    unit_type = 'slide'

    def __init__(self, raw: bytes, path: Path):
        super().__init__(raw, path)
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx required for PPTX support: pip install docrouter[pptx]")
        self._doc = Presentation(io.BytesIO(raw))

    def text(self) -> str:
        return '\n\n'.join(self.units())

    def units(self) -> list[str]:
        slides_text = []
        for slide in self._doc.slides:
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_parts.append(' | '.join(cells))
                elif hasattr(shape, 'text') and shape.text.strip():
                    slide_parts.append(shape.text)
            # Include speaker notes if available
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"[Notes: {notes}]")
            slides_text.append('\n'.join(slide_parts))
        return slides_text

    def meta(self) -> dict:
        props = self._doc.core_properties
        result = {}
        if props.title:
            result['title'] = props.title
        if props.author:
            result['author'] = props.author
        return result

# === HtmlDoc ===
class HtmlDoc(Doc):
    fmt = 'html'
    unit_type = 'section'

    def __init__(self, raw: bytes, path: Path):
        super().__init__(raw, path)
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4 required for HTML support: pip install docrouter[html]")
        self._soup = BeautifulSoup(raw, 'html.parser')
        # Remove script/style
        for tag in self._soup(['script', 'style']):
            tag.decompose()

    def text(self) -> str:
        return self._soup.get_text(separator='\n\n')

    def meta(self) -> dict:
        if self._soup.title and self._soup.title.string:
            return {'title': self._soup.title.string}
        return {}

# === ImgDoc ===
class ImgDoc(Doc):
    fmt = 'image'
    unit_type = 'none'

    def text(self) -> str:
        return ''

    def units(self) -> list[str]:
        return []

    def meta(self) -> dict:
        try:
            from PIL import Image
            img = Image.open(io.BytesIO(self._raw))
            return {'width': img.width, 'height': img.height, 'mode': img.mode}
        except Exception:
            return {}

# === Parser registry ===
_PARSERS: dict[str, type[Doc]] = {
    '.pdf': PDFDoc,
    '.docx': DocxDoc,
    '.pptx': PptxDoc,
    '.html': HtmlDoc, '.htm': HtmlDoc,
    '.md': MdDoc, '.markdown': MdDoc,
    '.txt': TxtDoc, '.text': TxtDoc, '.log': TxtDoc, '.csv': TxtDoc, '.json': TxtDoc,
    '.py': TxtDoc, '.js': TxtDoc, '.ts': TxtDoc, '.jsx': TxtDoc, '.tsx': TxtDoc,
    '.java': TxtDoc, '.c': TxtDoc, '.cpp': TxtDoc, '.h': TxtDoc, '.hpp': TxtDoc,
    '.rs': TxtDoc, '.go': TxtDoc, '.rb': TxtDoc, '.php': TxtDoc, '.swift': TxtDoc,
    '.kt': TxtDoc, '.scala': TxtDoc, '.sh': TxtDoc, '.bash': TxtDoc, '.zsh': TxtDoc,
    '.yaml': TxtDoc, '.yml': TxtDoc, '.toml': TxtDoc, '.ini': TxtDoc, '.cfg': TxtDoc,
    '.xml': TxtDoc, '.sql': TxtDoc, '.r': TxtDoc, '.m': TxtDoc, '.lua': TxtDoc,
    '.png': ImgDoc, '.jpg': ImgDoc, '.jpeg': ImgDoc, '.gif': ImgDoc, '.webp': ImgDoc, '.bmp': ImgDoc,
}

# === Search ===
def _search(doc: Doc, query: str, before: int, after: int, max_hits: int,
            mode: str, unit_start: int | None, unit_end: int | None, case_sensitive: bool) -> list[dict]:
    """Search for literal query in document units."""
    units = doc.units()
    if not units: return []

    # Validate range
    start = unit_start if unit_start is not None else 0
    end = unit_end if unit_end is not None else len(units) - 1
    if start < 0 or end >= len(units) or start > end:
        raise ValueError(f"invalid unit range [{start}, {end}] for document with {len(units)} units")

    results = []
    q = query if case_sensitive else query.lower()

    for ui in range(start, end + 1):
        unit_text = units[ui]
        search_text = unit_text if case_sensitive else unit_text.lower()
        pos = 0
        while len(results) < max_hits:
            idx = search_text.find(q, pos)
            if idx == -1: break

            if mode == 'unit':
                snippet = unit_text
            else:  # window
                s = max(0, idx - before)
                e = min(len(unit_text), idx + len(query) + after)
                snippet = unit_text[s:e]

            # Unit preview: first 200 chars
            preview = unit_text[:200] + ('...' if len(unit_text) > 200 else '')

            results.append({
                'unit_index': ui,
                'unit_type': doc.unit_type,
                'match_start': idx,
                'match_end': idx + len(query),
                'snippet': snippet,
                'unit_preview': preview
            })
            pos = idx + 1

        if len(results) >= max_hits: break

    return results

# === DocumentHandle ===
class DocumentHandle:
    """Main public interface for document operations."""

    def __init__(self, doc: Doc, path: Path, raw: bytes, doc_id: str):
        self._doc, self._path, self._raw, self._id = doc, path, raw, doc_id

    @property
    def document_id(self) -> str: return self._id

    def info(self) -> dict:
        """Return document metadata."""
        units = self._doc.units()
        has_text = bool(self._doc.text().strip())
        return {
            'document_id': self._id,
            'filename': self._path.name,
            'mime_type': _MIMES.get(self._doc.fmt, _MIMES['unknown']),
            'doc_type': self._doc.fmt,
            'unit_type': self._doc.unit_type,
            'unit_count': len(units),
            'has_text': has_text,
            **self._doc.meta()
        }

    def get_file_path(self) -> str:
        """Return path to original file."""
        return str(self._path)

    def get_bytes(self) -> bytes:
        """Return original file bytes."""
        return self._raw

    def get_text(self) -> str:
        """Return full document text."""
        return self._doc.text()

    def get_unit_text(self, unit_index: int) -> dict:
        """Return text for specific unit."""
        units = self._doc.units()
        if not 0 <= unit_index < len(units):
            raise IndexError(f"unit {unit_index} out of range (0-{len(units)-1})")
        return {
            'document_id': self._id,
            'unit_index': unit_index,
            'unit_type': self._doc.unit_type,
            'text': units[unit_index]
        }

    def search(self, query: str, *, before: int = 300, after: int = 300, max_hits: int = 5,
               mode: str = 'window', unit_start: int | None = None, unit_end: int | None = None,
               case_sensitive: bool = False) -> dict:
        """Search for literal substring. Returns hits ordered by unit_index, match_start."""
        if mode not in ('window', 'unit'):
            raise ValueError(f"mode must be 'window' or 'unit', got '{mode}'")
        hits = _search(self._doc, query, before, after, max_hits, mode, unit_start, unit_end, case_sensitive)
        return {
            'document_id': self._id,
            'query': query,
            'mode': mode,
            'max_hits': max_hits,
            'hits': hits
        }

    def render_page(self, page_index: int, *, dpi: int = 150, image_format: str = 'png') -> dict:
        """Render PDF page to image. Returns path to temp file. PDF only."""
        if self._doc.fmt != 'pdf':
            raise UnsupportedOperationError(f"render_page only supported for PDF, not {self._doc.fmt}")
        png = self._doc.render(page_index, dpi)
        tmp_dir = Path(tempfile.gettempdir()) / 'docrouter' / self._id
        tmp_dir.mkdir(parents=True, exist_ok=True)
        img_path = tmp_dir / f"page_{page_index}.{image_format}"
        img_path.write_bytes(png)
        return {
            'document_id': self._id,
            'page_index': page_index,
            'dpi': dpi,
            'image_format': image_format,
            'image_path': str(img_path)
        }

# === Main entry point ===
def open_document(inp: Union[str, Path, bytes, BinaryIO], filename: str | None = None) -> DocumentHandle:
    """Open a document from path, bytes, or file-like object."""
    path, raw = _resolve_input(inp, filename)
    ext = path.suffix.lower()
    if ext not in _PARSERS:
        raise UnsupportedFileTypeError(f"unsupported file type: {ext}")
    try:
        doc = _PARSERS[ext](raw, path)
    except ImportError:
        raise
    except Exception as e:
        raise ParseError(f"failed to parse {ext}: {e}") from e
    doc_id = f"doc_{uuid.uuid4().hex[:12]}"
    return DocumentHandle(doc, path, raw, doc_id)
