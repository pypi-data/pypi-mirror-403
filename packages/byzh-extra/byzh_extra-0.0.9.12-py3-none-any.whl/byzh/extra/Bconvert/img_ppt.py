from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import numpy as np
import os


def b_imgs2ppt(
        image_folder,
        output_ppt="output.pptx",
        sorted_key=lambda x: int(x.split(".")[0]),
        base_width_in=10  # 固定基准宽度（英寸）
):
    images = sorted([
        file for file in os.listdir(image_folder)
        if file.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif"))
    ], key=sorted_key)

    prs = Presentation()

    # -------- 先收集所有宽高比 --------
    aspect_ratios = []
    img_sizes = []
    for img_name in images:
        img_path = os.path.join(image_folder, img_name)
        with Image.open(img_path) as img:
            width_px, height_px = img.size
            aspect_ratio = height_px / width_px
            aspect_ratios.append(aspect_ratio)
            img_sizes.append((img_path, width_px, height_px, aspect_ratio))

    # -------- 用中位数确定全局页面大小 --------
    median_aspect_ratio = float(np.median(aspect_ratios))
    prs.slide_width = Inches(base_width_in)
    prs.slide_height = Inches(base_width_in * median_aspect_ratio)

    # -------- 逐张插入图片 --------
    for img_path, width_px, height_px, aspect_ratio in img_sizes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页
        try:
            # 计算该图在全局页面下的缩放
            img_height_in = base_width_in * aspect_ratio
            img_width_in = base_width_in

            # 如果图片太高 → 以页面高度为准缩放
            if img_height_in > prs.slide_height.inches:
                img_height_in = prs.slide_height.inches
                img_width_in = img_height_in / aspect_ratio

            # 居中摆放
            left = (prs.slide_width.inches - img_width_in) / 2
            top = (prs.slide_height.inches - img_height_in) / 2

            slide.shapes.add_picture(
                img_path,
                Inches(left), Inches(top),
                width=Inches(img_width_in),
                height=Inches(img_height_in)
            )
        except Exception as e:
            print(f"图片 {img_path} 添加失败: {e}")

    prs.save(output_ppt)
    print(f"PPT已保存到: {output_ppt}")
