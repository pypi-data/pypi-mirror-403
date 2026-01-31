from pptx.slide import Slide
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from ..image import ImageData
from ..parser import extract_image_variables
from ..tracer import Tracer
import re

"""
Procesa variables de tipo ImageData en slides de PPTX.
TODO: Refactor para otro tipo de archivos. El tracer y el parser deberian estar a nivel de templating
"""

class ImageProcessor:
    """Procesa variables de tipo ImageData en slides de PPTX"""
    
    @staticmethod
    def process_slide_images(slide: Slide, context: dict, tracer: Tracer):
        """
        Procesa todas las variables ImageData en una slide, reemplazando
        los placeholders {{variable}} con las imágenes reales.
        """
        shapes_to_remove = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Buscar variables ImageData en texto
                text_content = shape.text_frame.text
                image_vars = extract_image_variables(text_content, context)
                
                if image_vars:
                    # Si hay variables de imagen en este shape
                    for image_var in image_vars:
                        image_data = ImageProcessor._get_nested_value(context, image_var)
                        if isinstance(image_data, ImageData):
                            # Verificar si el shape contiene SOLO la variable de imagen
                            clean_text = text_content.strip()
                            var_pattern = r"""^\{\{\s*""" + re.escape(image_var) + r"""\s*\}\}$"""
                            
                            if re.match(var_pattern, clean_text):
                                # El shape contiene SOLO la imagen, reemplazarlo completamente
                                ImageProcessor._replace_shape_with_image(
                                    slide, shape, image_var, image_data, tracer
                                )
                                shapes_to_remove.append(shape)
                            else:
                                # El shape tiene texto mixto, insertar imagen cerca
                                ImageProcessor._insert_image_near_text(
                                    slide, shape, image_var, image_data, tracer
                                )
        
        # Remover shapes que fueron reemplazados por imágenes
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
    
    @staticmethod
    def _get_nested_value(context: dict, var_path: str):
        """
        Obtiene el valor de una variable anidada (ej: client.logo)
        """
        value = context
        for part in var_path.split('.'):
            value = value[part]
        return value
    
    @staticmethod
    def _replace_shape_with_image(slide: Slide, text_shape, image_var: str, image_data: ImageData, tracer: Tracer):
        """
        Reemplaza completamente un shape de texto con una imagen.
        """
        try:
            # Obtener posición y tamaño del shape de texto
            left = text_shape.left
            top = text_shape.top
            # width = text_shape.width
            # height = text_shape.height
            
            # Obtener los bytes de la imagen
            image_bytes = image_data.get_bytes()
            image_stream = BytesIO(image_bytes)
            image_stream.seek(0)
            
            # Calcular dimensiones finales
            final_width = Inches(image_data.width) if image_data.width else None
            final_height = Inches(image_data.height) if image_data.height else None

            # Si hay límites máximos, aplicarlos manteniendo aspect ratio
            if image_data.max_width or image_data.max_height:
                final_width, final_height = ImageProcessor._calculate_constrained_size(
                    final_width, final_height, image_data.max_width, image_data.max_height
                )
            
            # Añadir la imagen al slide
            picture = slide.shapes.add_picture(
                image_stream, left, top, final_width, final_height
            )
            
            tracer.log_variable(f"{image_var}", f"Imagen insertada ({final_width}x{final_height})")
            
        except Exception as e:
            tracer.log_error(f"Error insertando imagen '{image_var}': {e}")
    
    @staticmethod
    def _insert_image_near_text(slide: Slide, text_shape, image_var: str, image_data: ImageData, tracer: Tracer):
        """
        Inserta una imagen cerca de un shape de texto que contiene contenido mixto.
        """
        try:
            # Calcular posición para la imagen (debajo del texto)
            left = text_shape.left
            top = text_shape.top + text_shape.height + Inches(0.1)  # 0.1" de separación
            
            # Obtener los bytes de la imagen
            image_bytes = image_data.get_bytes()
            image_stream = BytesIO(image_bytes)
            
            # Usar dimensiones especificadas o tamaño por defecto
            final_width = Inches(image_data.width) if image_data.width else Inches(2.0)
            final_height = Inches(image_data.height) if image_data.height else Inches(1.5)
            
            # Si hay límites máximos, aplicarlos
            if image_data.max_width or image_data.max_height:
                final_width, final_height = ImageProcessor._calculate_constrained_size(
                    final_width, final_height, image_data.max_width, image_data.max_height
                )
            
            # Añadir la imagen al slide
            picture = slide.shapes.add_picture(
                image_stream, left, top, final_width, final_height
            )
            
            tracer.log_variable(f"{image_var}", f"Imagen insertada cerca del texto ({final_width}x{final_height})")
            
        except Exception as e:
            tracer.log_error(f"Error insertando imagen '{image_var}': {e}")
    
    @staticmethod
    def _calculate_constrained_size(width, height, max_width=None, max_height=None):
        """
        Calcula el tamaño final manteniendo aspect ratio dentro de los límites.
        """
        if max_width:
            max_width = Inches(max_width)
        if max_height:
            max_height = Inches(max_height)
            
        # Calcular factor de escala
        scale_x = max_width / width if max_width and width > max_width else 1
        scale_y = max_height / height if max_height and height > max_height else 1
        scale = min(scale_x, scale_y)
        
        return width * scale, height * scale