# -*- coding: utf-8 -*-
"""
Understand and analyze file contents using OpenAI's gpt-4.1 API.
Supports text files, PDF, DOCX, XLSX, and more.
"""

import json
import os
from pathlib import Path
from typing import List, Optional, Tuple

from dotenv import load_dotenv
from openai import AsyncOpenAI

from massgen.tool._result import ExecutionResult, TextContent


def _validate_path_access(path: Path, allowed_paths: Optional[List[Path]] = None) -> None:
    """
    Validate that a path is within allowed directories.

    Args:
        path: Path to validate
        allowed_paths: List of allowed base paths (optional)

    Raises:
        ValueError: If path is not within allowed directories
    """
    if not allowed_paths:
        return  # No restrictions

    for allowed_path in allowed_paths:
        try:
            path.relative_to(allowed_path)
            return  # Path is within this allowed directory
        except ValueError:
            continue

    raise ValueError(f"Path not in allowed directories: {path}")


def _extract_text_from_pdf(file_path: Path) -> Tuple[str, str]:
    """
    Extract text from a PDF file.

    Args:
        file_path: Path to the PDF file

    Returns:
        Tuple of (extracted_text, error_message)
        If successful, error_message is empty string
    """
    try:
        import PyPDF2
    except ImportError:
        return "", "PyPDF2 is required for PDF files. Install it with: pip install PyPDF2"

    try:
        text_content = []
        with open(file_path, "rb") as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            num_pages = len(pdf_reader.pages)

            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text.strip():
                    text_content.append(f"--- Page {page_num + 1} ---\n{text}")

        if not text_content:
            return "", "PDF file appears to be empty or contains only images"

        return "\n\n".join(text_content), ""

    except Exception as e:
        return "", f"Failed to extract text from PDF: {str(e)}"


def _extract_text_from_docx(file_path: Path) -> Tuple[str, str]:
    """
    Extract text from a DOCX file.

    Args:
        file_path: Path to the DOCX file

    Returns:
        Tuple of (extracted_text, error_message)
        If successful, error_message is empty string
    """
    try:
        from docx import Document
    except ImportError:
        return "", "python-docx is required for DOCX files. Install it with: pip install python-docx"

    try:
        doc = Document(file_path)
        text_content = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    text_content.append(row_text)

        if not text_content:
            return "", "DOCX file appears to be empty"

        return "\n\n".join(text_content), ""

    except Exception as e:
        return "", f"Failed to extract text from DOCX: {str(e)}"


def _extract_text_from_excel(file_path: Path) -> Tuple[str, str]:
    """
    Extract text from an Excel file (XLSX/XLS).

    Args:
        file_path: Path to the Excel file

    Returns:
        Tuple of (extracted_text, error_message)
        If successful, error_message is empty string
    """
    try:
        import openpyxl
    except ImportError:
        return "", "openpyxl is required for XLSX files. Install it with: pip install openpyxl"

    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text_content = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content.append(f"=== Sheet: {sheet_name} ===\n")

            for row in sheet.iter_rows(values_only=True):
                # Filter out None values and convert to string
                row_values = [str(cell) if cell is not None else "" for cell in row]
                # Only add non-empty rows
                if any(val.strip() for val in row_values):
                    text_content.append(" | ".join(row_values))

        if len(text_content) <= len(workbook.sheetnames):
            return "", "Excel file appears to be empty"

        return "\n".join(text_content), ""

    except Exception as e:
        return "", f"Failed to extract text from Excel: {str(e)}"


def _extract_text_from_pptx(file_path: Path) -> Tuple[str, str]:
    """
    Extract text from a PowerPoint file (PPTX).

    Args:
        file_path: Path to the PPTX file

    Returns:
        Tuple of (extracted_text, error_message)
        If successful, error_message is empty string
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "", "python-pptx is required for PPTX files. Install it with: pip install python-pptx"

    try:
        prs = Presentation(file_path)
        text_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"--- Slide {slide_num} ---")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)

        if len(text_content) <= len(prs.slides):
            return "", "PowerPoint file appears to be empty"

        return "\n\n".join(text_content), ""

    except Exception as e:
        return "", f"Failed to extract text from PowerPoint: {str(e)}"


async def understand_file(
    file_path: str,
    prompt: str = "Please analyze this file and provide a comprehensive understanding of its content, purpose, and structure.",
    model: str = "gpt-4.1",
    max_chars: int = 50000,
    allowed_paths: Optional[List[str]] = None,
    agent_cwd: Optional[str] = None,
    task_context: Optional[str] = None,
) -> ExecutionResult:
    """
    Understand and analyze file contents using OpenAI's gpt-4.1 API.

    This tool reads a file (text or document format) and processes its content through
    OpenAI's gpt-4.1 API to provide insights, summaries, explanations, or answer questions.

    Args:
        file_path: Path to the file to analyze
                  - Relative path: Resolved relative to workspace
                  - Absolute path: Must be within allowed directories
        prompt: Question or instruction about the file (default: asks for comprehensive analysis)
        model: Model to use (default: "gpt-4.1")
        max_chars: Maximum number of characters to read/extract (default: 50000)
                  - Prevents processing extremely large files
                  - Applies to both text files and extracted content from documents
        allowed_paths: List of allowed base paths for validation (optional)
        agent_cwd: Agent's current working directory (automatically injected, optional)

    Returns:
        ExecutionResult containing:
        - success: Whether operation succeeded
        - operation: "understand_file"
        - file_path: Path to the analyzed file
        - file_name: Name of the file
        - file_type: Extraction method used ("text", "pdf", "docx", "excel", "pptx")
        - file_size: Size of the file in bytes
        - chars_read: Number of characters read/extracted
        - truncated: Whether content was truncated
        - prompt: The prompt used
        - model: Model used for analysis
        - response: The model's understanding/analysis of the file

    Examples:
        # Text and code files
        understand_file("script.py")
        → Returns analysis of the Python script

        understand_file("README.md", "Summarize the key points of this documentation")
        → Returns summary of documentation

        # PDF documents
        understand_file("report.pdf", "What are the main findings in this research paper?")
        → Extracts text from PDF and analyzes it

        # Word documents
        understand_file("proposal.docx", "Summarize this business proposal")
        → Extracts text from DOCX and summarizes

        # Excel spreadsheets
        understand_file("data.xlsx", "What patterns can you see in this data?")
        → Extracts data from Excel and analyzes

        # PowerPoint presentations
        understand_file("presentation.pptx", "What are the key points of this presentation?")
        → Extracts text from slides and summarizes

    Security:
        - Requires valid OpenAI API key
        - File must exist and be readable
        - File content is sent to OpenAI API

    Supported File Types:
        Text Files:
        - Code: .py, .js, .java, .cpp, .c, .go, .rs, .ts, .tsx, .jsx, etc.
        - Config: .md, .yaml, .yml, .json, .xml, .toml, .ini, etc.
        - Data: .txt, .log, .csv, .tsv, etc.

        Document Files (require additional packages):
        - PDF: .pdf (requires PyPDF2: pip install PyPDF2)
        - Word: .docx (requires python-docx: pip install python-docx)
        - Excel: .xlsx (requires openpyxl: pip install openpyxl)
        - PowerPoint: .pptx (requires python-pptx: pip install python-pptx)

    Note:
        - Old Office formats (.doc, .xls, .ppt) are not supported
        - For images, use understand_image tool
        - For videos, use understand_video tool
        - For audio, use generate_text_with_input_audio tool
    """
    try:
        # Convert allowed_paths from strings to Path objects
        allowed_paths_list = [Path(p) for p in allowed_paths] if allowed_paths else None

        # Load environment variables
        script_dir = Path(__file__).parent.parent.parent.parent
        env_path = script_dir / ".env"
        if env_path.exists():
            load_dotenv(env_path)
        else:
            load_dotenv()

        openai_api_key = os.getenv("OPENAI_API_KEY")

        if not openai_api_key:
            result = {
                "success": False,
                "operation": "understand_file",
                "error": "OpenAI API key not found. Please set OPENAI_API_KEY in .env file or environment variable.",
            }
            return ExecutionResult(
                output_blocks=[TextContent(data=json.dumps(result, indent=2))],
            )

        # Initialize async OpenAI client
        client = AsyncOpenAI(api_key=openai_api_key)

        # Resolve file path
        # Use agent_cwd if available, otherwise fall back to Path.cwd()
        base_dir = Path(agent_cwd) if agent_cwd else Path.cwd()

        if Path(file_path).is_absolute():
            f_path = Path(file_path).resolve()
        else:
            f_path = (base_dir / file_path).resolve()

        # Validate file path
        _validate_path_access(f_path, allowed_paths_list)

        if not f_path.exists():
            result = {
                "success": False,
                "operation": "understand_file",
                "error": f"File does not exist: {f_path}",
            }
            return ExecutionResult(
                output_blocks=[TextContent(data=json.dumps(result, indent=2))],
            )

        if not f_path.is_file():
            result = {
                "success": False,
                "operation": "understand_file",
                "error": f"Path is not a file: {f_path}",
            }
            return ExecutionResult(
                output_blocks=[TextContent(data=json.dumps(result, indent=2))],
            )

        # Get file size
        file_size = f_path.stat().st_size

        # Check if file is unsupported binary format
        unsupported_binary_extensions = {
            ".exe",
            ".bin",
            ".dll",
            ".so",
            ".dylib",
            ".o",
            ".a",
            ".doc",  # Old Word format (use .docx instead)
            ".xls",  # Old Excel format (use .xlsx instead, though we try to support it)
            ".ppt",  # Old PowerPoint format (use .pptx instead)
            ".zip",
            ".tar",
            ".gz",
            ".bz2",
            ".7z",
            ".rar",
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".bmp",
            ".ico",
            ".svg",
            ".mp3",
            ".wav",
            ".ogg",
            ".flac",
            ".aac",
            ".mp4",
            ".avi",
            ".mov",
            ".mkv",
            ".flv",
            ".wmv",
            ".pyc",
            ".class",
            ".jar",
        }

        file_extension = f_path.suffix.lower()

        if file_extension in unsupported_binary_extensions:
            result = {
                "success": False,
                "operation": "understand_file",
                "error": f"Unsupported file format: {f_path.suffix}. " f"For images use understand_image, for videos use understand_video, for audio use generate_text_with_input_audio.",
            }
            return ExecutionResult(
                output_blocks=[TextContent(data=json.dumps(result, indent=2))],
            )

        # Extract content based on file type
        file_content = ""
        extraction_method = "text"

        # PDF files
        if file_extension == ".pdf":
            extraction_method = "pdf"
            file_content, error = _extract_text_from_pdf(f_path)
            if error:
                result = {
                    "success": False,
                    "operation": "understand_file",
                    "error": error,
                }
                return ExecutionResult(
                    output_blocks=[TextContent(data=json.dumps(result, indent=2))],
                )

        # Word documents
        elif file_extension == ".docx":
            extraction_method = "docx"
            file_content, error = _extract_text_from_docx(f_path)
            if error:
                result = {
                    "success": False,
                    "operation": "understand_file",
                    "error": error,
                }
                return ExecutionResult(
                    output_blocks=[TextContent(data=json.dumps(result, indent=2))],
                )

        # Excel spreadsheets
        elif file_extension in [".xlsx", ".xls"]:
            extraction_method = "excel"
            file_content, error = _extract_text_from_excel(f_path)
            if error:
                result = {
                    "success": False,
                    "operation": "understand_file",
                    "error": error,
                }
                return ExecutionResult(
                    output_blocks=[TextContent(data=json.dumps(result, indent=2))],
                )

        # PowerPoint presentations
        elif file_extension == ".pptx":
            extraction_method = "pptx"
            file_content, error = _extract_text_from_pptx(f_path)
            if error:
                result = {
                    "success": False,
                    "operation": "understand_file",
                    "error": error,
                }
                return ExecutionResult(
                    output_blocks=[TextContent(data=json.dumps(result, indent=2))],
                )

        # Text-based files
        else:
            try:
                with open(f_path, "r", encoding="utf-8") as file:
                    file_content = file.read(max_chars)

            except UnicodeDecodeError:
                # File is likely binary
                result = {
                    "success": False,
                    "operation": "understand_file",
                    "error": f"File appears to be binary and cannot be read as text: {f_path}. Supported binary formats: PDF, DOCX, XLSX, PPTX",
                }
                return ExecutionResult(
                    output_blocks=[TextContent(data=json.dumps(result, indent=2))],
                )
            except Exception as read_error:
                result = {
                    "success": False,
                    "operation": "understand_file",
                    "error": f"Failed to read file: {str(read_error)}",
                }
                return ExecutionResult(
                    output_blocks=[TextContent(data=json.dumps(result, indent=2))],
                )

        # Truncate if necessary (for text files)
        chars_read = len(file_content)
        truncated = False

        if extraction_method == "text" and chars_read == max_chars and file_size > max_chars:
            truncated = True
            truncation_note = f"\n\n[Note: File was truncated. Read {chars_read} characters out of {file_size} bytes total. Increase max_chars parameter to read more.]"
            file_content += truncation_note
        elif chars_read > max_chars:
            # Truncate extracted content from document formats
            truncated = True
            file_content = file_content[:max_chars]
            truncation_note = f"\n\n[Note: Extracted content was truncated. Showing first {max_chars} characters. Increase max_chars parameter to read more.]"
            file_content += truncation_note
            chars_read = len(file_content)

        # Inject task context into prompt if available
        from massgen.context.task_context import format_prompt_with_context

        augmented_prompt = format_prompt_with_context(prompt, task_context)

        # Build the full prompt with file content
        full_prompt = f"{augmented_prompt}\n\nFile: {f_path.name}\nContent:\n```\n{file_content}\n```"

        try:
            # Call OpenAI API for file understanding
            response = await client.responses.create(
                model=model,
                input=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "input_text", "text": full_prompt},
                        ],
                    },
                ],
            )

            # Extract response text
            response_text = response.output_text if hasattr(response, "output_text") else str(response.output)

            result = {
                "success": True,
                "operation": "understand_file",
                "file_path": str(f_path),
                "file_name": f_path.name,
                "file_type": extraction_method,
                "file_size": file_size,
                "chars_read": chars_read,
                "truncated": truncated,
                "prompt": prompt,
                "model": model,
                "response": response_text,
            }
            return ExecutionResult(
                output_blocks=[TextContent(data=json.dumps(result, indent=2))],
            )

        except Exception as api_error:
            result = {
                "success": False,
                "operation": "understand_file",
                "error": f"OpenAI API error: {str(api_error)}",
            }
            return ExecutionResult(
                output_blocks=[TextContent(data=json.dumps(result, indent=2))],
            )

    except Exception as e:
        result = {
            "success": False,
            "operation": "understand_file",
            "error": f"Failed to understand file: {str(e)}",
        }
        return ExecutionResult(
            output_blocks=[TextContent(data=json.dumps(result, indent=2))],
        )
