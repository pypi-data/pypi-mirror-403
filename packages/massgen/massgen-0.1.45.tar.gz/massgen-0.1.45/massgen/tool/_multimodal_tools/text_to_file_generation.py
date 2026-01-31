# -*- coding: utf-8 -*-
"""
Generate text content using OpenAI API and save it as various file formats (TXT, MD, PDF).
"""

import json
import os
from datetime import datetime
from pathlib import Path
from typing import List, Optional

from dotenv import load_dotenv
from openai import AsyncOpenAI

from massgen.tool._result import ExecutionResult, TextContent


def _validate_path_access(path: Path, allowed_paths: Optional[List[Path]] = None) -> None:
    """
    Validate that a path is within allowed directories.

    Args:
        path: Path to validate
        allowed_paths: List of allowed base paths (optional)
        agent_cwd: Agent\'s current working directory (automatically injected)

    Raises:
        ValueError: If path is not within allowed directories
    """
    if not allowed_paths:
        return  # No restrictions

    for allowed_path in allowed_paths:
        try:
            path.relative_to(allowed_path)
            return  # Path is within this allowed directory
        except ValueError:
            continue

    raise ValueError(f"Path not in allowed directories: {path}")


def _generate_pdf(content: str, file_path: Path) -> None:
    """
    Generate a PDF file from text content.

    Args:
        content: Text content to write to PDF
        file_path: Path where PDF will be saved
    """
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

        # Create PDF
        doc = SimpleDocTemplate(str(file_path), pagesize=letter)
        story = []
        styles = getSampleStyleSheet()

        # Split content into paragraphs
        paragraphs = content.split("\n\n")

        for para in paragraphs:
            if para.strip():
                # Handle special markdown-like formatting
                if para.startswith("#"):
                    # Use heading style
                    p = Paragraph(para.replace("#", "").strip(), styles["Heading1"])
                else:
                    p = Paragraph(para.replace("\n", "<br/>"), styles["BodyText"])
                story.append(p)
                story.append(Spacer(1, 0.2 * inch))

        doc.build(story)

    except ImportError:
        # Fallback: use fpdf if reportlab is not available
        try:
            from fpdf import FPDF

            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)

            # Split content into lines and add to PDF
            for line in content.split("\n"):
                pdf.multi_cell(0, 10, txt=line)

            pdf.output(str(file_path))

        except ImportError:
            raise ImportError(
                "PDF generation requires either 'reportlab' or 'fpdf2' library. " "Install with: pip install reportlab  OR  pip install fpdf2",
            )


def _generate_pptx(content: str, file_path: Path) -> None:
    """
    Generate a PowerPoint presentation from text content.

    Args:
        content: Text content to convert to PPTX (expects slide-based structure)
        file_path: Path where PPTX will be saved
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Parse content into slides
        # Expected format: slides separated by "---" or "Slide X:" markers
        # Or parse based on headers (##)

        slides_content = []
        current_slide = {"title": "", "content": []}

        lines = content.split("\n")
        i = 0

        while i < len(lines):
            line = lines[i].strip()

            # Check for slide delimiter
            if line.startswith("---") or line.startswith("==="):
                if current_slide["title"] or current_slide["content"]:
                    slides_content.append(current_slide)
                    current_slide = {"title": "", "content": []}
                i += 1
                continue

            # Check for title (marked with # or ##)
            if line.startswith("# "):
                if current_slide["title"] or current_slide["content"]:
                    slides_content.append(current_slide)
                    current_slide = {"title": "", "content": []}
                current_slide["title"] = line.lstrip("#").strip()
                i += 1
                continue

            # Check for subtitle/section (## or "Slide X:")
            if line.startswith("## ") or line.lower().startswith("slide "):
                if current_slide["title"] or current_slide["content"]:
                    slides_content.append(current_slide)
                    current_slide = {"title": "", "content": []}
                current_slide["title"] = line.lstrip("#").strip()
                i += 1
                continue

            # Add content to current slide
            if line:
                current_slide["content"].append(line)

            i += 1

        # Add last slide if it has content
        if current_slide["title"] or current_slide["content"]:
            slides_content.append(current_slide)

        # If no slides were parsed, create a single slide with all content
        if not slides_content:
            slides_content = [
                {
                    "title": "Generated Content",
                    "content": [line.strip() for line in content.split("\n") if line.strip()],
                },
            ]

        # Create slides
        for slide_data in slides_content:
            # Add title slide if it's the first slide and has only title
            if len(prs.slides) == 0 and slide_data["title"] and not slide_data["content"]:
                slide_layout = prs.slide_layouts[0]  # Title slide
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = slide_data["title"]
            else:
                # Add title and content slide
                slide_layout = prs.slide_layouts[1]  # Title and content
                slide = prs.slides.add_slide(slide_layout)

                # Set title
                title = slide.shapes.title
                title.text = slide_data["title"] if slide_data["title"] else "Content"

                # Set content
                if len(slide.shapes) > 1:
                    content_shape = slide.shapes[1]
                    text_frame = content_shape.text_frame
                    text_frame.clear()

                    for idx, content_line in enumerate(slide_data["content"]):
                        if idx == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()

                        # Handle bullet points
                        if content_line.startswith("- ") or content_line.startswith("* "):
                            p.text = content_line[2:].strip()
                            p.level = 0
                        elif content_line.startswith("  - ") or content_line.startswith("  * "):
                            p.text = content_line[4:].strip()
                            p.level = 1
                        else:
                            p.text = content_line
                            p.level = 0

        # Save presentation
        prs.save(str(file_path))

    except ImportError:
        raise ImportError(
            "PPTX generation requires 'python-pptx' library. " "Install with: pip install python-pptx",
        )


async def text_to_file_generation(
    prompt: str,
    file_format: str = "txt",
    filename: Optional[str] = None,
    model: str = "gpt-4o",
    storage_path: Optional[str] = None,
    allowed_paths: Optional[List[str]] = None,
    agent_cwd: Optional[str] = None,
) -> ExecutionResult:
    """
    Generate text content using OpenAI API and save it as various file formats.

    This tool uses OpenAI's chat completion API to generate text content based on a prompt,
    then saves the generated content in the specified file format (TXT, MD, PDF, or PPTX).

    Args:
        prompt: Description of the content to generate (e.g., "Write a technical report about AI")
        file_format: Output file format - Options: "txt", "md", "pdf", "pptx" (default: "txt")
        filename: Custom filename without extension (optional)
                 If not provided, generates from prompt and timestamp
        model: OpenAI model to use (default: "gpt-4o")
               Options: "gpt-4o", "gpt-4o-mini", "gpt-4-turbo", "gpt-3.5-turbo"
        storage_path: Directory path where to save the file (optional)
                     - **IMPORTANT**: Must be a DIRECTORY path only, NOT a file path (e.g., "documents/reports" NOT "documents/report.txt")
                     - The filename is automatically generated from the prompt or custom filename parameter
                     - Relative path: Resolved relative to agent's workspace (e.g., "documents/reports")
                     - Absolute path: Must be within allowed directories
                     - None/empty: Saves to agent's workspace root
        allowed_paths: List of allowed base paths for validation (optional)
        agent_cwd: Agent\'s current working directory (automatically injected)

    Returns:
        ExecutionResult containing:
        - success: Whether operation succeeded
        - operation: "generate_and_store_file"
        - file_path: Path to the generated file
        - filename: Name of the generated file
        - file_format: Format of the generated file
        - content_preview: First 500 characters of generated content
        - file_size: Size of the generated file in bytes
        - model: Model used for generation
        - prompt: The prompt used

    Examples:
        text_to_file_generation("Write a blog post about Python", file_format="md")
        → Generates markdown file with blog post content

        text_to_file_generation(
            "Create a technical report on machine learning",
            file_format="pdf",
            filename="ml_report"
        )
        → Generates PDF file named "ml_report.pdf"

        text_to_file_generation(
            "Write meeting notes for today's standup",
            file_format="txt",
            storage_path="documents/notes"
        )
        → Generates text file in documents/notes/ directory

    Security:
        - Requires valid OpenAI API key
        - Files are saved to specified path within workspace
        - Path must be within allowed directories

    Note:
        - PDF generation requires either 'reportlab' or 'fpdf2' library
        - PPTX generation requires 'python-pptx' library
        - For PPTX format, structure your prompt to include slide titles (using # or ##) and bullet points (using -)
        - The quality and format of generated content depends on the prompt
        - Longer content may consume more tokens
    """
    try:
        # Validate file format
        supported_formats = ["txt", "md", "pdf", "pptx"]
        file_format = file_format.lower()
        if file_format not in supported_formats:
            result = {
                "success": False,
                "operation": "generate_and_store_file",
                "error": f"Unsupported file format: {file_format}. Supported formats: {', '.join(supported_formats)}",
            }
            return ExecutionResult(
                output_blocks=[TextContent(data=json.dumps(result, indent=2))],
            )

        # Convert allowed_paths from strings to Path objects
        allowed_paths_list = [Path(p) for p in allowed_paths] if allowed_paths else None

        # Use agent_cwd if available, otherwise fall back to base_dir
        base_dir = Path(agent_cwd) if agent_cwd else Path.cwd()

        # Load environment variables
        script_dir = Path(__file__).parent.parent.parent.parent
        env_path = script_dir / ".env"
        if env_path.exists():
            load_dotenv(env_path)
        else:
            load_dotenv()

        openai_api_key = os.getenv("OPENAI_API_KEY")

        if not openai_api_key:
            result = {
                "success": False,
                "operation": "generate_and_store_file",
                "error": "OpenAI API key not found. Please set OPENAI_API_KEY in .env file or environment variable.",
            }
            return ExecutionResult(
                output_blocks=[TextContent(data=json.dumps(result, indent=2))],
            )

        # Initialize async OpenAI client
        client = AsyncOpenAI(api_key=openai_api_key)

        # Determine storage directory
        if storage_path:
            if Path(storage_path).is_absolute():
                storage_dir = Path(storage_path).resolve()
            else:
                storage_dir = (base_dir / storage_path).resolve()
        else:
            storage_dir = base_dir

        # Validate storage directory
        _validate_path_access(storage_dir, allowed_paths_list)
        storage_dir.mkdir(parents=True, exist_ok=True)

        try:
            # Generate content using OpenAI API
            response = await client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "system",
                        "content": f"You are a professional content writer. Generate high-quality {file_format.upper()} content based on the user's request.",
                    },
                    {
                        "role": "user",
                        "content": prompt,
                    },
                ],
                temperature=0.7,
            )

            # Extract generated content
            generated_content = response.choices[0].message.content

            if not generated_content:
                result = {
                    "success": False,
                    "operation": "generate_and_store_file",
                    "error": "No content generated from OpenAI API",
                }
                return ExecutionResult(
                    output_blocks=[TextContent(data=json.dumps(result, indent=2))],
                )

        except Exception as api_error:
            result = {
                "success": False,
                "operation": "generate_and_store_file",
                "error": f"OpenAI API error: {str(api_error)}",
            }
            return ExecutionResult(
                output_blocks=[TextContent(data=json.dumps(result, indent=2))],
            )

        # Generate filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        if filename:
            # Use custom filename (remove extension if provided)
            clean_filename = filename.rsplit(".", 1)[0]
            file_name = f"{clean_filename}.{file_format}"
        else:
            # Generate filename from prompt
            clean_prompt = "".join(c for c in prompt[:30] if c.isalnum() or c in (" ", "-", "_")).strip()
            clean_prompt = clean_prompt.replace(" ", "_")
            file_name = f"{timestamp}_{clean_prompt}.{file_format}"

        # Full file path
        file_path = storage_dir / file_name

        # Save content based on format
        try:
            if file_format == "pdf":
                _generate_pdf(generated_content, file_path)
            elif file_format == "pptx":
                _generate_pptx(generated_content, file_path)
            else:
                # For txt and md, save as plain text
                file_path.write_text(generated_content, encoding="utf-8")

            file_size = file_path.stat().st_size

        except Exception as save_error:
            result = {
                "success": False,
                "operation": "generate_and_store_file",
                "error": f"Failed to save file: {str(save_error)}",
            }
            return ExecutionResult(
                output_blocks=[TextContent(data=json.dumps(result, indent=2))],
            )

        # Create result
        result = {
            "success": True,
            "operation": "generate_and_store_file",
            "file_path": str(file_path),
            "filename": file_name,
            "file_format": file_format,
            "content_preview": generated_content[:500] + ("..." if len(generated_content) > 500 else ""),
            "file_size": file_size,
            "model": model,
            "prompt": prompt,
        }

        return ExecutionResult(
            output_blocks=[TextContent(data=json.dumps(result, indent=2))],
        )

    except Exception as e:
        result = {
            "success": False,
            "operation": "generate_and_store_file",
            "error": f"Failed to generate file: {str(e)}",
        }
        return ExecutionResult(
            output_blocks=[TextContent(data=json.dumps(result, indent=2))],
        )
