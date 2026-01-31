"""filetextio.parsers

Text ingestion utilities for multiple file formats.

Provides a simple API to load text from PDF, DOCX, PPTX, and TXT files.
"""

from __future__ import annotations

from abc import ABC, abstractmethod
from pathlib import Path
from typing import Dict, Type, Union
from llm_lab.exceptions import UnsupportedFileTypeError, ParseError
import base64
import io
import mimetypes
import os

# Optional heavy dependencies (safe import for new projects)
try:
    import pypdfium2 as pdfium  # type: ignore[import-not-found]
except Exception:
    pdfium = None  # type: ignore[assignment]

try:
    from pdf2image import convert_from_path  # type: ignore[import-not-found]
except Exception:
    convert_from_path = None  # type: ignore[assignment]

try:
    import pdfplumber  # type: ignore[import-not-found]
except Exception:
    pdfplumber = None  # type: ignore[assignment]

PathLike = Union[str, Path]



class FileParser(ABC):
    """Abstract base for file parsers returning plain text.

    Subclasses must implement `read_file()` to return a best-effort textual
    representation of the file content.
    """

    def __init__(self, path: PathLike):
        """Initialize a file parser.

        Args:
            path: Path to the target file to be parsed. Accepts a string path or
                a ``pathlib.Path``.
        """
        self.path = str(path)

    @abstractmethod
    def read_file(self) -> str:
        """Read the file and return extracted text.

        Raises:
            ParseError: if the file cannot be parsed.
        """
        raise NotImplementedError


class PDFParser(FileParser):
    def read_file(self) -> str:
        """Extract text from a PDF using pdfplumber.

        Returns:
            Concatenated text from all pages.

        Raises:
            ParseError: if pdfplumber fails to open or parse the file.
        """
        if pdfplumber is None:
            raise ParseError("pdfplumber is not installed; cannot parse PDFs.")
        try:
            parts: list[str] = []
            with pdfplumber.open(self.path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    parts.append(text)
            return "\n".join(parts)
        except Exception as exc:
            raise ParseError(f"Failed to parse PDF: {self.path}") from exc

    @classmethod
    def pdf_to_base64_images(cls, pdf_path: str, max_pages: int | None = None) -> list[str]:
        """
        Convert PDF pages to base64-encoded JPEG images for vision workflows.

        Args:
            pdf_path: Path to the PDF file on disk
            max_pages: Maximum number of pages to convert.

        Returns:
            List of base64-encoded JPEG image strings.
        """
        try:
            images = []

            # Prefer rendering with PDFium when available
            if pdfium is not None:
                doc = pdfium.PdfDocument(pdf_path)
                try:
                    page_count = len(doc)
                    limit = page_count if max_pages is None else min(page_count, max_pages)
                    for i in range(limit):
                        # Access page (supports both __getitem__ and get_page APIs)
                        try:
                            page = doc.get_page(i)  # type: ignore[attr-defined]
                        except AttributeError:
                            page = doc[i]
                        # Render at 2x scale (~144 DPI) for clarity
                        bitmap = page.render(scale=2.0)
                        img = bitmap.to_pil()
                        images.append(img)
                        # Cleanup page resources if API provides close()
                        try:
                            page.close()  # type: ignore[attr-defined]
                        except Exception:
                            pass
                finally:
                    try:
                        doc.close()
                    except Exception:
                        pass

            # Fallback to pdf2image if PDFium unavailable or produced no images
            if not images and convert_from_path is not None:
                images = convert_from_path(pdf_path, dpi=150, fmt="jpeg")
                if max_pages is not None:
                    images = images[:max_pages]

            # Convert to base64
            base64_images = []
            for img in images:
                # Convert to RGB if needed
                if img.mode != "RGB":
                    img = img.convert("RGB")

                # Save to bytes buffer
                buffer = io.BytesIO()
                img.save(buffer, format="JPEG", quality=85)
                buffer.seek(0)

                # Encode to base64
                img_base64 = base64.b64encode(buffer.read()).decode("utf-8")
                base64_images.append(img_base64)

            return base64_images

        except Exception:
            return []



class DocxParser(FileParser):
    def read_file(self) -> str:
        """Extract text from a DOCX file using python-docx.

        Returns:
            Concatenated text from paragraph runs.

        Raises:
            ParseError: if the DOCX cannot be opened or parsed.
        """
        from docx import Document

        try:
            doc = Document(self.path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as exc:
            raise ParseError(f"Failed to parse DOCX: {self.path}") from exc


class PPTXParser(FileParser):
    def read_file(self) -> str:
        """Extract text from a PPTX file using python-pptx.

        Returns:
            Concatenated text grouped by slides.

        Raises:
            ParseError: if the PPTX cannot be opened or parsed.
        """
        from pptx import Presentation

        try:
            prs = Presentation(self.path)
            parts: list[str] = []

            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_text: list[str] = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text_frame:
                        text = "\n".join(
                            p.text for p in shape.text_frame.paragraphs if p.text
                        )
                        if text.strip():
                            slide_text.append(text.strip())
                    elif hasattr(shape, "text"):
                        text = (shape.text or "").strip()
                        if text:
                            slide_text.append(text)

                if slide_text:
                    parts.append(f"[SLIDE {slide_idx}]\n" + "\n".join(slide_text))

            return "\n\n".join(parts)
        except Exception as exc:
            raise ParseError(f"Failed to parse PPTX: {self.path}") from exc


class TXTParser(FileParser):
    def read_file(self) -> str:
        """Read plain text files safely with UTF-8 and ignore errors."""
        try:
            return Path(self.path).read_text(encoding="utf-8", errors="ignore")
        except Exception as exc:
            raise ParseError(f"Failed to read TXT: {self.path}") from exc


 


_PARSER_REGISTRY: Dict[str, Type[FileParser]] = {
    ".pdf": PDFParser,
    ".docx": DocxParser,
    ".pptx": PPTXParser,
    ".txt": TXTParser,
}

# Supported extension set for quick checks
_SUPPORTED_EXTENSIONS = set(_PARSER_REGISTRY.keys())


def detect_file_format(path: PathLike, prefer_mime: bool = True) -> str:
    """Best-effort detection of a file's format.

    Tries MIME type via ``mimetypes.guess_type`` first (when ``prefer_mime``),
    then falls back to the file suffix. Returns a normalized extension like
    ".pdf", ".docx", ".pptx", or ".txt". Returns an empty string if unknown.

    Args:
        path: File path to inspect.
        prefer_mime: Prefer MIME-based detection before suffix.

    Returns:
        Normalized extension including the leading dot, or "" if unknown.
    """
    p = Path(path)
    ext = p.suffix.lower()

    if prefer_mime:
        mime, _ = mimetypes.guess_type(str(p))
        if mime == "application/pdf":
            return ".pdf"
        if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return ".docx"
        if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return ".pptx"
        if mime and mime.startswith("text/"):
            return ".txt"

    # Fall back to suffix if MIME didn't help
    return ext if ext in _SUPPORTED_EXTENSIONS else ""


def get_parser_for(path: PathLike) -> FileParser:
    """Return an appropriate parser instance based on detected file format.

    Args:
        path: Path to a file on disk.

    Raises:
        FileNotFoundError: If the path does not exist.
        IsADirectoryError: If the path points to a directory.
        UnsupportedFileTypeError: If the format is not supported.
    """
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")
    if p.is_dir():
        raise IsADirectoryError(f"Path is a directory, expected file: {p}")

    ext = detect_file_format(p)
    cls = _PARSER_REGISTRY.get(ext)
    if not cls:
        raise UnsupportedFileTypeError(f"Unsupported file type: {ext or 'unknown'}")
    return cls(str(p))


def load_text(path: PathLike) -> str:
    """Convenience function: parse and return text for any supported file.

    Args:
        path: Path to a file on disk.

    Raises:
        FileNotFoundError: If the path does not exist.
        IsADirectoryError: If the path points to a directory.
        UnsupportedFileTypeError: If the extension is not supported.
        ParseError: If the parser fails to extract text.
    """
    return get_parser_for(path).read_file()


def get_supported_types() -> list[str]:
    """Return the list of supported normalized file extensions."""
    return sorted(_SUPPORTED_EXTENSIONS)


def is_supported_file(path: PathLike) -> bool:
    """Check whether the given file path appears to be supported.

    Uses the same detection strategy as ``get_parser_for`` but returns a bool.
    """
    try:
        p = Path(path)
        if not p.exists() or p.is_dir():
            return False
        return detect_file_format(p, prefer_mime=True) in _SUPPORTED_EXTENSIONS
    except Exception:
        return False


# Utility helpers for temporary file handling in async upload flows
async def create_temp_pdf_file(file, temp_path, logger) -> bool:
    """Write an uploaded PDF to a temporary path.

    Args:
        file: Async file-like object supporting ``await file.read()`` and ``filename``.
        temp_path: Destination path on disk for the temporary file.
        logger: Logger with a ``warning`` method.

    Returns:
        True if content was written, False if the upload was empty.
    """
    with open(temp_path, "wb") as f:
        content = await file.read()
        if not content:
            try:
                filename = getattr(file, "filename", "<unknown>")
            except Exception:
                filename = "<unknown>"
            logger.warning(f"Empty file uploaded: {filename}")
            return False
        f.write(content)
    return True


def delete_file(temp_path) -> None:
    """Delete a file if it exists, ignoring errors."""
    if os.path.exists(temp_path):
        try:
            os.remove(temp_path)
        except Exception:
            pass


__all__ = [
    "FileParser",
    "PDFParser",
    "DocxParser",
    "PPTXParser",
    "TXTParser",
    "UnsupportedFileTypeError",
    "ParseError",
    "detect_file_format",
    "get_supported_types",
    "is_supported_file",
    "get_parser_for",
    "load_text",
    "create_temp_pdf_file",
    "delete_file",
]
