import json
import os
from typing import List, Dict, Any, Optional
from ..utils.helpers import logger

def load_txt_file(file_path: str) -> str:
    """Load text from a file"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        logger.error(f"Failed to load text file: {str(e)}")
        raise

def load_json_file(file_path: str) -> Any:
    """Load JSON from a file"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        logger.error(f"Failed to load JSON file: {str(e)}")
        raise

def load_docx_file(file_path: str) -> str:
    """Load text from DOCX file"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + "\n"
        
        return text.strip()
    except ImportError:
        raise ImportError("python-docx is required for DOCX support. Install with: pip install kssrag[office]")
    except Exception as e:
        logger.error(f"Failed to load DOCX file: {str(e)}")
        raise

def load_excel_file(file_path: str) -> str:
    """Load text from Excel file"""
    try:
        import openpyxl
        workbook = openpyxl.load_workbook(file_path)
        text = ""
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"Sheet: {sheet_name}\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text += row_text + "\n"
            text += "\n"
        
        return text.strip()
    except ImportError:
        raise ImportError("openpyxl is required for Excel support. Install with: pip install kssrag[office]")
    except Exception as e:
        logger.error(f"Failed to load Excel file: {str(e)}")
        raise

def load_pptx_file(file_path: str) -> str:
    """Load text from PowerPoint file"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        
        for slide_number, slide in enumerate(prs.slides, 1):
            text += f"Slide {slide_number}:\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            
            text += "\n"
        
        return text.strip()
    except ImportError:
        raise ImportError("python-pptx is required for PowerPoint support. Install with: pip install kssrag[office]")
    except Exception as e:
        logger.error(f"Failed to load PowerPoint file: {str(e)}")
        raise

def load_document(file_path: str) -> str:
    """Load document from file with auto-format detection"""
    if file_path.endswith('.txt'):
        return load_txt_file(file_path)
    elif file_path.endswith('.docx'):
        return load_docx_file(file_path)
    elif file_path.endswith(('.xlsx', '.xls')):
        return load_excel_file(file_path)
    elif file_path.endswith('.pptx'):
        return load_pptx_file(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_path}")

def load_json_documents(file_path: str, metadata_field: str = "name") -> List[Dict[str, Any]]:
    """Load documents from JSON file"""
    data = load_json_file(file_path)
    
    # Apply limit for testing if specified
    from ..config import config
    if config.MAX_DOCS_FOR_TESTING:
        data = data[:config.MAX_DOCS_FOR_TESTING]
        logger.info(f"Limited to {config.MAX_DOCS_FOR_TESTING} documents for testing")
    
    return data