"""
Document Preprocessor - Converts various document formats to text.

Supports:
- Plain text files (.txt, .md, .py, .js, .json, etc.)
- Microsoft Word (.docx)
- PDF files (.pdf)
- Excel spreadsheets (.xlsx, .xls)
- PowerPoint presentations (.pptx)

The preprocessor creates cached text versions of binary documents
for efficient RLM access.
"""

import hashlib
import json
import logging
import os
import re
import warnings
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Set, Tuple

from watsonx_rlm_knowledge.exceptions import PreprocessingError

logger = logging.getLogger(__name__)


# File extensions that are already text-readable
TEXT_EXTENSIONS: Set[str] = {
    # Code files
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cc", ".cpp", ".h", ".hpp",
    ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala", ".cs", ".fs", ".clj",
    ".hs", ".ml", ".r", ".jl", ".lua", ".pl", ".pm", ".sh", ".bash", ".zsh", ".fish",
    ".ps1", ".psm1", ".bat", ".cmd", ".awk", ".sed",
    # Config files
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".xml", ".xsl", ".xslt", ".dtd", ".xsd", ".plist",
    # Documentation
    ".md", ".markdown", ".rst", ".txt", ".text", ".asciidoc", ".adoc",
    ".org", ".wiki", ".tex", ".latex",
    # Web files
    ".html", ".htm", ".xhtml", ".css", ".scss", ".sass", ".less",
    ".vue", ".svelte", ".astro",
    # Data files
    ".csv", ".tsv", ".sql", ".graphql", ".gql",
    # Other
    ".proto", ".thrift", ".avsc", ".gradle", ".properties",
    ".gitignore", ".dockerignore", ".editorconfig",
    "Makefile", "Dockerfile", "Vagrantfile", "Jenkinsfile",
}

# Extensions that need conversion
BINARY_DOC_EXTENSIONS: Set[str] = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    ".odt", ".ods", ".odp",  # OpenDocument formats
    ".rtf",
}


@dataclass
class PreprocessedDocument:
    """Metadata about a preprocessed document."""
    original_path: str
    text_path: str
    original_hash: str
    format: str
    size_bytes: int
    text_size_bytes: int
    page_count: Optional[int] = None
    sheet_names: Optional[List[str]] = None
    slide_count: Optional[int] = None


@dataclass
class PreprocessorConfig:
    """Configuration for document preprocessing."""
    cache_dir: str = ".rlm_cache"
    max_file_size_mb: int = 50
    skip_hidden: bool = True
    skip_dirs: Tuple[str, ...] = (
        ".git", ".hg", ".svn", "node_modules", "__pycache__",
        ".venv", "venv", ".env", "dist", "build", ".tox",
        ".pytest_cache", ".mypy_cache", ".ruff_cache",
    )
    include_extensions: Optional[Set[str]] = None  # None = all supported
    exclude_patterns: Tuple[str, ...] = (
        "*.min.js", "*.min.css", "*.map", "*.lock",
        "package-lock.json", "yarn.lock", "poetry.lock",
    )


class DocumentPreprocessor:
    """Preprocesses documents for RLM knowledge access.

    Converts binary document formats (PDF, DOCX, XLSX, PPTX) to text files
    that can be efficiently searched and read by the RLM engine.

    Text files are cached with content hashes to avoid re-processing
    unchanged documents.

    Example:
        preprocessor = DocumentPreprocessor("/path/to/knowledge")
        preprocessor.preprocess_all()

        # Get text content of a document
        text = preprocessor.get_text("/path/to/knowledge/report.pdf")
    """

    def __init__(self, root: str, config: Optional[PreprocessorConfig] = None):
        """Initialize preprocessor.

        Args:
            root: Root directory of knowledge base
            config: Optional preprocessing configuration
        """
        self.root = Path(root).resolve()
        self.config = config or PreprocessorConfig()
        self.cache_dir = self.root / self.config.cache_dir
        self.manifest_path = self.cache_dir / "manifest.json"
        self._manifest: Dict[str, PreprocessedDocument] = {}

        # Ensure root exists
        if not self.root.exists():
            raise PreprocessingError(f"Knowledge root does not exist: {self.root}")

        # Load existing manifest
        self._load_manifest()

    def _load_manifest(self):
        """Load preprocessing manifest from cache."""
        if self.manifest_path.exists():
            try:
                with open(self.manifest_path, "r") as f:
                    data = json.load(f)
                    for path, doc_data in data.items():
                        self._manifest[path] = PreprocessedDocument(**doc_data)
                logger.info(f"Loaded manifest with {len(self._manifest)} documents")
            except Exception as e:
                logger.warning(f"Failed to load manifest: {e}")
                self._manifest = {}

    def _save_manifest(self):
        """Save preprocessing manifest to cache."""
        self.cache_dir.mkdir(parents=True, exist_ok=True)

        data = {
            path: {
                "original_path": doc.original_path,
                "text_path": doc.text_path,
                "original_hash": doc.original_hash,
                "format": doc.format,
                "size_bytes": doc.size_bytes,
                "text_size_bytes": doc.text_size_bytes,
                "page_count": doc.page_count,
                "sheet_names": doc.sheet_names,
                "slide_count": doc.slide_count,
            }
            for path, doc in self._manifest.items()
        }

        with open(self.manifest_path, "w") as f:
            json.dump(data, f, indent=2)

    def _compute_hash(self, path: Path) -> str:
        """Compute SHA256 hash of file content."""
        hasher = hashlib.sha256()
        with open(path, "rb") as f:
            while chunk := f.read(65536):
                hasher.update(chunk)
        return hasher.hexdigest()[:16]

    def _should_skip_path(self, path: Path) -> bool:
        """Check if path should be skipped."""
        # Skip hidden files/dirs
        if self.config.skip_hidden:
            for part in path.relative_to(self.root).parts:
                if part.startswith(".") and part not in {".env"}:
                    return True

        # Skip excluded directories
        for part in path.parts:
            if part in self.config.skip_dirs:
                return True

        # Skip excluded patterns
        for pattern in self.config.exclude_patterns:
            if path.match(pattern):
                return True

        return False

    def _is_supported(self, path: Path) -> bool:
        """Check if file is supported for processing."""
        suffix = path.suffix.lower()

        # Check explicit include list
        if self.config.include_extensions:
            return suffix in self.config.include_extensions

        # Check known extensions
        return suffix in TEXT_EXTENSIONS or suffix in BINARY_DOC_EXTENSIONS

    def _is_binary_doc(self, path: Path) -> bool:
        """Check if file needs conversion from binary."""
        return path.suffix.lower() in BINARY_DOC_EXTENSIONS

    def _needs_processing(self, path: Path) -> bool:
        """Check if file needs (re)processing."""
        path_str = str(path)

        if path_str not in self._manifest:
            return True

        doc = self._manifest[path_str]

        # Check if original file changed
        current_hash = self._compute_hash(path)
        if current_hash != doc.original_hash:
            return True

        # Check if text file exists
        if not Path(doc.text_path).exists():
            return True

        return False

    def iter_documents(self) -> List[Path]:
        """Iterate over all documents in the knowledge base."""
        documents = []

        for path in self.root.rglob("*"):
            if path.is_file():
                if self._should_skip_path(path):
                    continue
                if self._is_supported(path):
                    # Check file size
                    size_mb = path.stat().st_size / (1024 * 1024)
                    if size_mb <= self.config.max_file_size_mb:
                        documents.append(path)
                    else:
                        logger.warning(f"Skipping large file: {path} ({size_mb:.1f}MB)")

        return sorted(documents)

    def preprocess_all(self, force: bool = False) -> Dict[str, PreprocessedDocument]:
        """Preprocess all documents in the knowledge base.

        Args:
            force: If True, reprocess all documents regardless of cache

        Returns:
            Dictionary mapping original paths to preprocessed document info
        """
        documents = self.iter_documents()
        logger.info(f"Found {len(documents)} documents to process")

        processed = 0
        skipped = 0
        errors = 0

        for path in documents:
            try:
                if force or self._needs_processing(path):
                    self._preprocess_document(path)
                    processed += 1
                else:
                    skipped += 1
            except Exception as e:
                logger.error(f"Failed to process {path}: {e}")
                errors += 1

        # Save updated manifest
        self._save_manifest()

        logger.info(
            f"Preprocessing complete: {processed} processed, "
            f"{skipped} skipped, {errors} errors"
        )

        return self._manifest.copy()

    def _preprocess_document(self, path: Path) -> PreprocessedDocument:
        """Preprocess a single document."""
        path_str = str(path)
        suffix = path.suffix.lower()
        file_hash = self._compute_hash(path)

        if self._is_binary_doc(path):
            # Convert binary document to text
            text, metadata = self._convert_to_text(path)

            # Store in cache
            text_filename = f"{path.stem}_{file_hash}.txt"
            text_path = self.cache_dir / "texts" / text_filename
            text_path.parent.mkdir(parents=True, exist_ok=True)

            with open(text_path, "w", encoding="utf-8") as f:
                f.write(text)

            doc = PreprocessedDocument(
                original_path=path_str,
                text_path=str(text_path),
                original_hash=file_hash,
                format=suffix,
                size_bytes=path.stat().st_size,
                text_size_bytes=len(text.encode("utf-8")),
                **metadata
            )
        else:
            # Text file - reference directly
            doc = PreprocessedDocument(
                original_path=path_str,
                text_path=path_str,  # Same as original
                original_hash=file_hash,
                format=suffix,
                size_bytes=path.stat().st_size,
                text_size_bytes=path.stat().st_size,
            )

        self._manifest[path_str] = doc
        logger.debug(f"Processed: {path}")

        return doc

    def _convert_to_text(self, path: Path) -> Tuple[str, Dict]:
        """Convert binary document to text.

        Returns:
            Tuple of (text_content, metadata_dict)
        """
        suffix = path.suffix.lower()

        if suffix == ".pdf":
            return self._convert_pdf(path)
        elif suffix in (".docx", ".doc"):
            return self._convert_docx(path)
        elif suffix in (".xlsx", ".xls"):
            return self._convert_xlsx(path)
        elif suffix in (".pptx", ".ppt"):
            return self._convert_pptx(path)
        elif suffix == ".rtf":
            return self._convert_rtf(path)
        else:
            raise PreprocessingError(f"Unsupported format: {suffix}")

    def _convert_pdf(self, path: Path) -> Tuple[str, Dict]:
        """Convert PDF to text."""
        try:
            from pypdf import PdfReader
        except ImportError:
            raise PreprocessingError("pypdf not installed. Run: pip install pypdf")

        import sys
        import io
        import logging as _logging

        # Suppress pypdf warnings about malformed PDFs (common and harmless)
        # These include "Ignoring wrong pointing object" messages printed to stderr
        pypdf_logger = _logging.getLogger("pypdf")
        old_level = pypdf_logger.level
        pypdf_logger.setLevel(_logging.ERROR)

        # Capture stderr to suppress pypdf's direct print statements
        old_stderr = sys.stderr
        sys.stderr = io.StringIO()

        try:
            with warnings.catch_warnings():
                warnings.filterwarnings("ignore", message=".*wrong pointing object.*")
                warnings.filterwarnings("ignore", message=".*Xref.*")
                warnings.filterwarnings("ignore", module="pypdf")

                reader = PdfReader(str(path))
                pages = []

                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append(f"--- Page {i + 1} ---\n{text}")
        finally:
            sys.stderr = old_stderr
            pypdf_logger.setLevel(old_level)

        return "\n\n".join(pages), {"page_count": len(reader.pages)}

    def _convert_docx(self, path: Path) -> Tuple[str, Dict]:
        """Convert DOCX to text."""
        try:
            from docx import Document
        except ImportError:
            raise PreprocessingError("python-docx not installed. Run: pip install python-docx")

        doc = Document(str(path))
        paragraphs = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Preserve heading structure
                if para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "")
                    try:
                        level_num = int(level)
                        prefix = "#" * level_num + " "
                    except ValueError:
                        prefix = "# "
                    paragraphs.append(prefix + text)
                else:
                    paragraphs.append(text)

        # Also extract tables
        tables_text = []
        for i, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                tables_text.append(f"\n--- Table {i + 1} ---\n" + "\n".join(rows))

        full_text = "\n\n".join(paragraphs)
        if tables_text:
            full_text += "\n" + "\n".join(tables_text)

        return full_text, {"page_count": None}

    def _convert_xlsx(self, path: Path) -> Tuple[str, Dict]:
        """Convert XLSX to text."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise PreprocessingError("openpyxl not installed. Run: pip install openpyxl")

        wb = load_workbook(str(path), read_only=True, data_only=True)
        sheets = []
        sheet_names = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_names.append(sheet_name)
            rows = []

            for row in sheet.iter_rows(values_only=True):
                # Convert row to strings, handling None
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in cells):  # Skip empty rows
                    rows.append(" | ".join(cells))

            if rows:
                sheets.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))

        wb.close()
        return "\n\n".join(sheets), {"sheet_names": sheet_names}

    def _convert_pptx(self, path: Path) -> Tuple[str, Dict]:
        """Convert PPTX to text."""
        try:
            from pptx import Presentation
        except ImportError:
            raise PreprocessingError("python-pptx not installed. Run: pip install python-pptx")

        prs = Presentation(str(path))
        slides = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"--- Slide {i + 1} ---")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # Handle tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_text.append(" | ".join(cells))

            if len(slide_text) > 1:  # More than just the header
                slides.append("\n".join(slide_text))

        return "\n\n".join(slides), {"slide_count": len(prs.slides)}

    def _convert_rtf(self, path: Path) -> Tuple[str, Dict]:
        """Convert RTF to text (basic conversion)."""
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        # Basic RTF stripping (not perfect but handles simple RTFs)
        # Remove RTF control words
        text = re.sub(r"\\[a-z]+\d* ?", "", content)
        # Remove braces
        text = re.sub(r"[{}]", "", text)
        # Clean up whitespace
        text = re.sub(r"\s+", " ", text).strip()

        return text, {}

    def get_text(self, path: str) -> str:
        """Get text content of a document.

        Args:
            path: Path to original document

        Returns:
            Text content (converted if necessary)
        """
        path = str(Path(path).resolve())

        # Check manifest
        if path in self._manifest:
            doc = self._manifest[path]
            text_path = doc.text_path
        else:
            # Process on demand
            doc = self._preprocess_document(Path(path))
            text_path = doc.text_path

        with open(text_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    def get_document_info(self, path: str) -> Optional[PreprocessedDocument]:
        """Get preprocessing info for a document."""
        path = str(Path(path).resolve())
        return self._manifest.get(path)

    def get_all_text_paths(self) -> Dict[str, str]:
        """Get mapping of original paths to text paths."""
        return {
            doc.original_path: doc.text_path
            for doc in self._manifest.values()
        }