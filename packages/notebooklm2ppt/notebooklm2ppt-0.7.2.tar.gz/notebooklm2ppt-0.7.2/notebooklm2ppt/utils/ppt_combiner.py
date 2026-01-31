# 使用 Spire.Presentation 合并PPT文件，保留原始设计
# 安装: pip install spire.presentation

import os
from pathlib import Path
from spire.presentation import *
from spire.presentation.common import *
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation as PptxPresentation
from PIL import Image

def combine_ppt_files_with_spire(source_folder, output_file, png_names=None):
    """
    使用 Spire.Presentation 合并PPT文件，每个PPT只保留第一页，并保留原始设计
    
    Args:
        source_folder: 源PPT文件所在的文件夹路径
        output_file: 输出的合并PPT文件路径
    """
    # 获取所有pptx文件并按字典序排序
    
    ppt_files = sorted([f for f in os.listdir(source_folder) if f.endswith('.pptx')])

    ppt_names = [png_name.replace(".png", ".pptx") for png_name in png_names] if png_names else None
    
    if not ppt_files:
        print("未找到任何PPT文件")
        return
    
    print(f"找到 {len(ppt_files)} 个PPT文件:")
    valid_ppt_files = []
    for idx, file in enumerate(ppt_files, 1):
        print(f"  {idx}. {file}")
        if ppt_names and file in ppt_names:
            valid_ppt_files.append(file)
    if ppt_names:
        ppt_files = valid_ppt_files
        print(f"\n根据提供的PNG名称过滤后，剩余 {len(ppt_files)} 个PPT文件:")
    
    # 创建主演示文稿对象，使用第一个PPT作为基础
    first_ppt_path = os.path.join(source_folder, ppt_files[0])
    main_pres = Presentation()
    main_pres.LoadFromFile(first_ppt_path)
    
    # 删除第一个PPT的多余页面，只保留第一页
    while main_pres.Slides.Count > 1:
        main_pres.Slides.RemoveAt(1)
    
    print(f"  已添加: {ppt_files[0]} (第1页)")
    
    # 遍历剩余的PPT文件
    for ppt_file in ppt_files[1:]:
        file_path = os.path.join(source_folder, ppt_file)
        
        # 加载当前PPT文件
        temp_pres = Presentation()
        temp_pres.LoadFromFile(file_path)
        
        if temp_pres.Slides.Count > 0:
            # 使用 AppendBySlide 方法追加第一页，保留原始设计
            main_pres.Slides.AppendBySlide(temp_pres.Slides[0])
            print(f"  已添加: {ppt_file} (第1页)")
        else:
            print(f"  跳过: {ppt_file} (无幻灯片)")
        
        # 释放临时演示文稿资源
        temp_pres.Dispose()
    
    # 保存合并后的PPT
    main_pres.SaveToFile(output_file, FileFormat.Pptx2016)
    print(f"\n合并完成！输出文件: {output_file}")
    print(f"总共合并了 {main_pres.Slides.Count} 页幻灯片")
    
    # 释放资源
    main_pres.Dispose()

    valid_png_names = [ppt_name.replace(".pptx", ".png") for ppt_name in ppt_files] # 返回真正存在的png names

    return valid_png_names



def clean_ppt(in_ppt_file, out_ppt_file):
    """
    使用 python-pptx 删除PPT中的多余形状
    
    Args:
        ppt_file: 需要清理的PPT文件路径
    """
    ppt = PptxPresentation(in_ppt_file)
    for slide in ppt.slides:
        for shape in list(slide.shapes):
            # 如果shape name 叫做"New shape",删除它
            if shape.name == "New shape":
                sp = slide.shapes._spTree.remove(shape._element)

    ppt.save(out_ppt_file)
    

def combine_ppt(source_folder, out_ppt_file, png_names = None):
    # 确保是字符串路径，因为后面用到了 .replace
    source_folder = str(source_folder)
    out_ppt_file = str(out_ppt_file)
    
    # 方法1: 保留原始设计（推荐）
    output_file1 = out_ppt_file.replace(".pptx", "_combined_original_design.pptx")
    valid_png_names = combine_ppt_files_with_spire(source_folder, output_file1, png_names=png_names)

    clean_ppt(output_file1, out_ppt_file)
    print(f"\n已生成合并的PPT文件: {out_ppt_file}")
    os.remove(output_file1)
    
    return valid_png_names


def create_ppt_from_images(png_dir, output_file, png_names=None):
    """
    直接将PNG图片插入到PPTX中，不使用智能圈选
    
    Args:
        png_dir: PNG图片所在目录
        output_file: 输出PPT文件路径
        png_names: PNG文件名列表（可选，如果提供则只处理这些文件）
    
    Returns:
        实际使用的PNG文件名列表
    """
    png_dir = Path(png_dir)
    output_file = str(output_file)
    
    if png_names:
        png_files = [png_dir / name for name in png_names if (png_dir / name).exists()]
    else:
        png_files = sorted(png_dir.glob("*.png"))
    
    if not png_files:
        print("未找到任何PNG图片")
        return []
    
    print(f"找到 {len(png_files)} 张PNG图片")
    
    prs = PptxPresentation()
    
    first_img = Image.open(png_files[0])
    img_width_px, img_height_px = first_img.size
    
    prs.slide_width = Pt(img_width_px)
    prs.slide_height = Pt(img_height_px)
    
    print(f"PPT尺寸设置为: {img_width_px} x {img_height_px} 像素")
    
    blank_layout = prs.slide_layouts[6]
    
    for idx, png_file in enumerate(png_files, 1):
        print(f"  [{idx}/{len(png_files)}] 处理: {png_file.name}")
        
        slide = prs.slides.add_slide(blank_layout)
        
        img = Image.open(png_file)
        img_width, img_height = img.size
        
        slide.shapes.add_picture(str(png_file), 0, 0, Pt(img_width), Pt(img_height))
    
    prs.save(output_file)
    print(f"\n已生成PPT文件: {output_file}")
    print(f"总共添加了 {len(prs.slides)} 页幻灯片")
    
    return [f.name for f in png_files]
